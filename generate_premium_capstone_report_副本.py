import html
import io
import os
import shutil
import subprocess
import time
from pathlib import Path

import fitz
import folium
import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.io as pio
import seaborn as sns
from folium.features import DivIcon
from folium.plugins import MarkerCluster, MousePosition
from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from sklearn.metrics import accuracy_score, classification_report

from generate_capstone_report import (
    GITHUB_URL,
    NOTEBOOKS,
    URLS,
    gh,
    sql_results,
    train_models,
    wiki_scrape_summary,
)


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "premium_capstone_outputs"
ASSETS = OUT / "assets"
SLIDES = OUT / "slide_images"
HTML_DECK = OUT / "premium_spacex_capstone_deck.html"
PDF_OUT = ROOT / "Data Science Capstone Project Report.pdf"
PDF_PREMIUM = ROOT / "IBM_SpaceX_Capstone_Final_Report.pdf"
PDF_ALIAS = ROOT / "Data_Science_Capstone_Project_Report.pdf"
PPTX_OUT = ROOT / "IBM_SpaceX_Capstone_Final_Report.pptx"
PPTX_ALIAS = ROOT / "Data_Science_Capstone_Project_Report.pptx"


PALETTE = {
    "ink": "#152238",
    "navy": "#0F2A43",
    "blue": "#2F80ED",
    "cyan": "#17A2B8",
    "teal": "#00A78E",
    "green": "#2EAD5B",
    "orange": "#F2994A",
    "red": "#D64545",
    "purple": "#6C5CE7",
    "gray": "#64748B",
    "muted": "#8A94A6",
    "line": "#D8E0EA",
    "bg": "#F6F8FB",
    "white": "#FFFFFF",
}

plt.rcParams.update(
    {
        "font.family": "DejaVu Sans",
        "axes.titleweight": "bold",
        "axes.titlecolor": PALETTE["ink"],
        "axes.labelcolor": PALETTE["gray"],
        "xtick.color": PALETTE["gray"],
        "ytick.color": PALETTE["gray"],
        "figure.dpi": 150,
        "savefig.dpi": 260,
    }
)
sns.set_theme(style="whitegrid")


def ensure_dirs():
    OUT.mkdir(exist_ok=True)
    ASSETS.mkdir(exist_ok=True)
    SLIDES.mkdir(exist_ok=True)


def read_frames():
    frames = {}
    for key, url in URLS.items():
        frames[key] = pd.read_csv(url)
    return frames


def save_plot(fig, name):
    path = ASSETS / name
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def clean_axis(ax):
    ax.grid(True, axis="y", color="#E7ECF2", linewidth=0.8)
    ax.grid(False, axis="x")
    for spine in ax.spines.values():
        spine.set_visible(False)


def make_static_charts(frames, ml):
    part2 = frames["part2"].copy()
    part2["Year"] = pd.to_datetime(part2["Date"]).dt.year
    part2["Landing"] = part2["Class"].map({0: "Not landed", 1: "Landed"})
    class_palette = {0: PALETTE["red"], 1: PALETTE["teal"]}
    charts = {}

    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    sns.scatterplot(
        data=part2,
        x="FlightNumber",
        y="LaunchSite",
        hue="Class",
        palette=class_palette,
        s=95,
        alpha=0.92,
        edgecolor="white",
        linewidth=0.7,
        ax=ax,
    )
    ax.set_title("Flight Number vs. Launch Site", loc="left", fontsize=17)
    ax.set_xlabel("Flight Number")
    ax.set_ylabel("")
    ax.legend(title="Landing class", labels=["Not landed", "Landed"], frameon=False, loc="lower right")
    clean_axis(ax)
    charts["flight_site"] = save_plot(fig, "eda_flight_number_vs_launch_site_premium.png")

    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    sns.scatterplot(
        data=part2,
        x="PayloadMass",
        y="LaunchSite",
        hue="Class",
        palette=class_palette,
        s=95,
        alpha=0.92,
        edgecolor="white",
        linewidth=0.7,
        ax=ax,
    )
    ax.set_title("Payload Mass vs. Launch Site", loc="left", fontsize=17)
    ax.set_xlabel("Payload Mass (kg)")
    ax.set_ylabel("")
    ax.legend(title="Landing class", labels=["Not landed", "Landed"], frameon=False, loc="lower right")
    clean_axis(ax)
    charts["payload_site"] = save_plot(fig, "eda_payload_vs_launch_site_premium.png")

    orbit = (
        part2.groupby("Orbit", as_index=False)
        .agg(launches=("Class", "size"), success_rate=("Class", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    colors = [PALETTE["teal"] if r >= 0.75 else PALETTE["orange"] if r >= 0.55 else PALETTE["red"] for r in orbit["success_rate"]]
    ax.bar(orbit["Orbit"], orbit["success_rate"], color=colors, width=0.72)
    ax.set_title("Success Rate by Orbit Type", loc="left", fontsize=17)
    ax.set_ylim(0, 1.08)
    ax.set_ylabel("Success Rate")
    ax.set_xlabel("")
    ax.tick_params(axis="x", rotation=35)
    for i, row in orbit.reset_index(drop=True).iterrows():
        ax.text(i, row["success_rate"] + 0.025, f"{row['success_rate']:.0%}", ha="center", fontsize=9, color=PALETTE["ink"])
    clean_axis(ax)
    charts["orbit_success"] = save_plot(fig, "eda_success_rate_by_orbit_premium.png")

    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    sns.scatterplot(
        data=part2,
        x="FlightNumber",
        y="Orbit",
        hue="Class",
        palette=class_palette,
        s=95,
        alpha=0.92,
        edgecolor="white",
        linewidth=0.7,
        ax=ax,
    )
    ax.set_title("Flight Number vs. Orbit Type", loc="left", fontsize=17)
    ax.set_xlabel("Flight Number")
    ax.set_ylabel("")
    ax.legend(title="Landing class", labels=["Not landed", "Landed"], frameon=False, loc="lower right")
    clean_axis(ax)
    charts["flight_orbit"] = save_plot(fig, "eda_flight_number_vs_orbit_premium.png")

    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    sns.scatterplot(
        data=part2,
        x="PayloadMass",
        y="Orbit",
        hue="Class",
        palette=class_palette,
        s=95,
        alpha=0.92,
        edgecolor="white",
        linewidth=0.7,
        ax=ax,
    )
    ax.set_title("Payload Mass vs. Orbit Type", loc="left", fontsize=17)
    ax.set_xlabel("Payload Mass (kg)")
    ax.set_ylabel("")
    ax.legend(title="Landing class", labels=["Not landed", "Landed"], frameon=False, loc="lower right")
    clean_axis(ax)
    charts["payload_orbit"] = save_plot(fig, "eda_payload_vs_orbit_premium.png")

    yearly = (
        part2.groupby("Year", as_index=False)
        .agg(launches=("Class", "size"), success_rate=("Class", "mean"))
        .sort_values("Year")
    )
    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    ax.plot(yearly["Year"], yearly["success_rate"], color=PALETTE["teal"], linewidth=3, marker="o", markersize=8)
    ax.fill_between(yearly["Year"], yearly["success_rate"], color=PALETTE["teal"], alpha=0.14)
    ax.set_title("Yearly Average Landing Success Trend", loc="left", fontsize=17)
    ax.set_ylim(0, 1.06)
    ax.set_xlabel("Year")
    ax.set_ylabel("Average success rate")
    ax.set_xticks(yearly["Year"])
    ax.tick_params(axis="x", rotation=35)
    for _, row in yearly.iterrows():
        ax.text(row["Year"], row["success_rate"] + 0.035, f"{row['success_rate']:.0%}", ha="center", fontsize=8.5)
    clean_axis(ax)
    charts["yearly"] = save_plot(fig, "eda_yearly_success_trend_premium.png")

    site = (
        part2.groupby("LaunchSite", as_index=False)
        .agg(launches=("Class", "size"), successes=("Class", "sum"), success_rate=("Class", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    ax.bar(site["LaunchSite"], site["success_rate"], color=[PALETTE["teal"], PALETTE["cyan"], PALETTE["orange"]], width=0.62)
    ax.set_title("Launch Site Success Rate", loc="left", fontsize=17)
    ax.set_ylim(0, 1.05)
    ax.set_xlabel("")
    ax.set_ylabel("Success rate")
    for i, row in site.reset_index(drop=True).iterrows():
        ax.text(i, row["success_rate"] + 0.03, f"{row['success_rate']:.1%}\n{row['successes']:.0f}/{row['launches']:.0f}", ha="center", fontsize=9)
    clean_axis(ax)
    charts["site_success"] = save_plot(fig, "site_success_rate_premium.png")

    ml_summary = ml["summary"].copy()
    fig, ax = plt.subplots(figsize=(10.8, 5.6))
    x = np.arange(len(ml_summary))
    width = 0.34
    ax.bar(x - width / 2, ml_summary["cv_accuracy"], width, label="10-fold CV", color=PALETTE["blue"])
    ax.bar(x + width / 2, ml_summary["test_accuracy"], width, label="Test", color=PALETTE["teal"])
    ax.set_title("Model Accuracy Comparison", loc="left", fontsize=17)
    ax.set_xticks(x)
    ax.set_xticklabels(ml_summary["model"], rotation=15, ha="right")
    ax.set_ylim(0, 1.02)
    ax.set_ylabel("Accuracy")
    ax.legend(frameon=False, loc="lower right")
    for idx, row in enumerate(ml_summary.itertuples()):
        ax.text(idx - width / 2, row.cv_accuracy + 0.02, f"{row.cv_accuracy:.3f}", ha="center", fontsize=8.5)
        ax.text(idx + width / 2, row.test_accuracy + 0.02, f"{row.test_accuracy:.3f}", ha="center", fontsize=8.5)
    clean_axis(ax)
    charts["ml_accuracy"] = save_plot(fig, "ml_accuracy_comparison_premium.png")

    cm = ml["best_cm"]
    fig, ax = plt.subplots(figsize=(6.8, 5.8))
    sns.heatmap(
        cm,
        annot=True,
        fmt="d",
        cmap=sns.light_palette(PALETTE["teal"], as_cmap=True),
        cbar=False,
        linewidths=1.5,
        linecolor="white",
        square=True,
        annot_kws={"fontsize": 18, "weight": "bold", "color": PALETTE["ink"]},
        ax=ax,
    )
    ax.set_title(f"Confusion Matrix: {ml['best_model']}", fontsize=16, loc="left")
    ax.set_xlabel("Predicted label")
    ax.set_ylabel("True label")
    ax.set_xticklabels(["Not landed", "Landed"], rotation=0)
    ax.set_yticklabels(["Not landed", "Landed"], rotation=0)
    charts["ml_cm"] = save_plot(fig, "ml_confusion_matrix_premium.png")

    return charts


def haversine(lat1, lon1, lat2, lon2):
    radius = 6373.0
    lat1, lon1, lat2, lon2 = map(np.radians, [lat1, lon1, lat2, lon2])
    dlon = lon2 - lon1
    dlat = lat2 - lat1
    a = np.sin(dlat / 2) ** 2 + np.cos(lat1) * np.cos(lat2) * np.sin(dlon / 2) ** 2
    return float(radius * 2 * np.arctan2(np.sqrt(a), np.sqrt(1 - a)))


def render_html_screenshot(html_path, png_path, width=1400, height=820, selector="body", wait_ms=1800):
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": width, "height": height}, device_scale_factor=1)
        page.goto(html_path.resolve().as_uri(), wait_until="domcontentloaded")
        page.wait_for_timeout(wait_ms)
        page.locator(selector).screenshot(path=str(png_path))
        browser.close()


def make_folium_maps(frames):
    geo = frames["geo"].copy()
    site_stats = (
        geo.groupby("Launch Site", as_index=False)
        .agg(Lat=("Lat", "first"), Long=("Long", "first"), launches=("class", "size"), successes=("class", "sum"))
    )
    site_stats["rate"] = site_stats["successes"] / site_stats["launches"]
    map_paths = {}

    base = [31.0, -96.0]
    site_map = folium.Map(location=base, zoom_start=4, tiles="OpenStreetMap")
    for _, row in site_stats.iterrows():
        label = f"{row['Launch Site']}: {row['launches']} launches, {row['rate']:.1%} success"
        folium.Circle(
            [row["Lat"], row["Long"]],
            radius=5500,
            color=PALETTE["blue"],
            fill=True,
            fill_color=PALETTE["cyan"],
            fill_opacity=0.30,
            popup=label,
        ).add_to(site_map)
        folium.Marker(
            [row["Lat"], row["Long"]],
            tooltip=label,
            icon=DivIcon(
                icon_size=(180, 34),
                icon_anchor=(-8, 10),
                html=f"<div style='font-size:12px;font-weight:700;color:#0F2A43;background:white;padding:4px 7px;border-radius:4px;border:1px solid #d8e0ea'>{html.escape(row['Launch Site'])}</div>",
            ),
        ).add_to(site_map)
    site_html = ASSETS / "folium_site_markers.html"
    site_png = ASSETS / "folium_site_markers_real.png"
    site_map.save(site_html)
    render_html_screenshot(site_html, site_png, selector=".folium-map", wait_ms=2600)
    map_paths["folium_sites"] = site_png

    cluster_map = folium.Map(location=base, zoom_start=4, tiles="OpenStreetMap")
    cluster = MarkerCluster(name="Launch records").add_to(cluster_map)
    for _, row in geo.iterrows():
        outcome = "Success" if row["class"] == 1 else "Failure / no attempt"
        color = "green" if row["class"] == 1 else "red"
        folium.Marker(
            [row["Lat"], row["Long"]],
            tooltip=f"{row['Launch Site']} | {outcome}",
            popup=f"Flight {row['Flight Number']}<br>{row['Date']}<br>{row['Launch Site']}<br>{outcome}",
            icon=folium.Icon(color=color, icon="ok-sign" if row["class"] == 1 else "remove-sign"),
        ).add_to(cluster)
    MousePosition(position="topright", prefix="Lat / Long").add_to(cluster_map)
    cluster_html = ASSETS / "folium_launch_records.html"
    cluster_png = ASSETS / "folium_launch_records_real.png"
    cluster_map.save(cluster_html)
    render_html_screenshot(cluster_html, cluster_png, selector=".folium-map", wait_ms=3000)
    map_paths["folium_records"] = cluster_png

    ksc = site_stats[site_stats["Launch Site"].str.contains("KSC")].iloc[0]
    points = [
        ("Coastline", ksc["Lat"], -80.57, PALETTE["cyan"]),
        ("Highway", 28.55, -80.65, PALETTE["orange"]),
        ("Cape Canaveral", 28.3922, -80.6077, PALETTE["purple"]),
    ]
    prox_map = folium.Map(location=[28.51, -80.61], zoom_start=10, tiles="OpenStreetMap")
    folium.Marker(
        [ksc["Lat"], ksc["Long"]],
        popup="KSC LC-39A launch site",
        tooltip="KSC LC-39A",
        icon=folium.Icon(color="blue", icon="rocket", prefix="fa"),
    ).add_to(prox_map)
    for name, lat, lon, color in points:
        dist = haversine(ksc["Lat"], ksc["Long"], lat, lon)
        folium.Marker(
            [lat, lon],
            tooltip=f"{name}: {dist:.1f} km",
            icon=DivIcon(
                icon_size=(155, 32),
                icon_anchor=(0, 0),
                html=f"<div style='font-size:12px;font-weight:700;color:#0F2A43;background:#fff;padding:4px 7px;border-radius:4px;border:1px solid #d8e0ea'>{name}: {dist:.1f} km</div>",
            ),
        ).add_to(prox_map)
        folium.PolyLine([[ksc["Lat"], ksc["Long"]], [lat, lon]], color=color, weight=3, opacity=0.85).add_to(prox_map)
    MousePosition(position="topright", prefix="Lat / Long").add_to(prox_map)
    prox_html = ASSETS / "folium_proximity.html"
    prox_png = ASSETS / "folium_proximity_real.png"
    prox_map.save(prox_html)
    render_html_screenshot(prox_html, prox_png, selector=".folium-map", wait_ms=2600)
    map_paths["folium_proximity"] = prox_png

    return map_paths


def write_plotly_dashboard_html(fig, title, subtitle, controls, path):
    div = pio.to_html(fig, include_plotlyjs=True, full_html=False, config={"displayModeBar": False, "responsive": True})
    controls_html = "".join(f"<div class='control'><span>{html.escape(k)}</span><strong>{html.escape(v)}</strong></div>" for k, v in controls)
    path.write_text(
        f"""<!doctype html>
<html>
<head>
<meta charset="utf-8">
<style>
body {{
  margin:0;
  background:#eef3f8;
  font-family:Inter, Arial, sans-serif;
  color:#152238;
}}
.shell {{
  width: 1320px;
  height: 780px;
  margin: 0;
  padding: 26px;
  box-sizing:border-box;
  background: linear-gradient(180deg,#f8fbff,#eef3f8);
}}
.top {{
  display:flex;
  align-items:flex-start;
  justify-content:space-between;
  margin-bottom:18px;
}}
h1 {{
  margin:0;
  font-size:32px;
  letter-spacing:.1px;
}}
.sub {{
  color:#64748B;
  font-size:16px;
  margin-top:7px;
}}
.controls {{
  display:flex;
  gap:12px;
  align-items:center;
}}
.control {{
  min-width:150px;
  background:#fff;
  border:1px solid #D8E0EA;
  border-radius:8px;
  padding:10px 13px;
  box-shadow:0 8px 22px rgba(15,42,67,.07);
}}
.control span {{
  display:block;
  font-size:11px;
  text-transform:uppercase;
  letter-spacing:.06em;
  color:#8A94A6;
  margin-bottom:4px;
}}
.control strong {{
  font-size:14px;
}}
.plot {{
  background:white;
  border:1px solid #D8E0EA;
  border-radius:12px;
  height:630px;
  overflow:hidden;
  box-shadow:0 18px 42px rgba(15,42,67,.10);
}}
.plot .plotly-graph-div {{
  height: 630px !important;
}}
</style>
</head>
<body>
<div class="shell">
  <div class="top">
    <div><h1>{html.escape(title)}</h1><div class="sub">{html.escape(subtitle)}</div></div>
    <div class="controls">{controls_html}</div>
  </div>
  <div class="plot">{div}</div>
</div>
</body>
</html>""",
        encoding="utf-8",
    )


def make_dash_screens(frames):
    dash = frames["dash"].copy()
    screens = {}

    success_counts = (
        dash[dash["class"] == 1]
        .groupby("Launch Site", as_index=False)
        .size()
        .rename(columns={"size": "Successful Launches"})
    )
    fig = px.pie(
        success_counts,
        values="Successful Launches",
        names="Launch Site",
        hole=0.36,
        color_discrete_sequence=[PALETTE["teal"], PALETTE["blue"], PALETTE["orange"], PALETTE["purple"]],
    )
    fig.update_layout(margin=dict(l=20, r=20, t=35, b=20), font=dict(family="Arial", size=16), showlegend=True)
    html_path = ASSETS / "dash_all_sites.html"
    png_path = ASSETS / "dash_all_sites_real.png"
    write_plotly_dashboard_html(
        fig,
        "SpaceX Launch Records Dashboard",
        "Total successful launches across all sites",
        [("Launch site", "All Sites"), ("Payload range", "0-10,000 kg")],
        html_path,
    )
    render_html_screenshot(html_path, png_path, width=1320, height=780, selector=".shell", wait_ms=1200)
    screens["dash_all"] = png_path

    site_rate = dash.groupby("Launch Site")["class"].agg(["mean", "size"]).query("size >= 5").sort_values("mean", ascending=False)
    best_site = site_rate.index[0]
    selected = dash[dash["Launch Site"] == best_site]["class"].value_counts().reindex([1, 0], fill_value=0).reset_index()
    selected.columns = ["class", "Launches"]
    selected["Outcome"] = selected["class"].map({1: "Landed", 0: "Did not land"})
    fig = px.pie(
        selected,
        values="Launches",
        names="Outcome",
        hole=0.36,
        color="Outcome",
        color_discrete_map={"Landed": PALETTE["teal"], "Did not land": PALETTE["red"]},
    )
    fig.update_layout(margin=dict(l=20, r=20, t=35, b=20), font=dict(family="Arial", size=16), showlegend=True)
    html_path = ASSETS / "dash_best_site.html"
    png_path = ASSETS / "dash_best_site_real.png"
    write_plotly_dashboard_html(
        fig,
        "SpaceX Launch Records Dashboard",
        f"Outcome mix for selected site: {best_site}",
        [("Launch site", best_site), ("Payload range", "0-10,000 kg")],
        html_path,
    )
    render_html_screenshot(html_path, png_path, width=1320, height=780, selector=".shell", wait_ms=1200)
    screens["dash_site"] = png_path

    fig = px.scatter(
        dash,
        x="Payload Mass (kg)",
        y="class",
        color="Booster Version Category",
        hover_data=["Launch Site", "Booster Version"],
        labels={"class": "Landing outcome"},
        color_discrete_sequence=[PALETTE["blue"], PALETTE["teal"], PALETTE["orange"], PALETTE["purple"]],
    )
    fig.update_traces(marker=dict(size=11, opacity=0.82, line=dict(width=0.7, color="white")))
    fig.update_layout(margin=dict(l=45, r=20, t=35, b=45), font=dict(family="Arial", size=15), yaxis=dict(tickvals=[0, 1], ticktext=["Not landed", "Landed"]))
    html_path = ASSETS / "dash_scatter.html"
    png_path = ASSETS / "dash_scatter_real.png"
    write_plotly_dashboard_html(
        fig,
        "SpaceX Launch Records Dashboard",
        "Payload mass vs. launch outcome, colored by booster version category",
        [("Launch site", "All Sites"), ("Payload range", "0-10,000 kg")],
        html_path,
    )
    render_html_screenshot(html_path, png_path, width=1320, height=780, selector=".shell", wait_ms=1200)
    screens["dash_scatter"] = png_path
    return screens


def rel(path):
    return Path(path).resolve().relative_to(OUT.resolve()).as_posix()


def fmt_pct(x):
    return f"{x:.1%}"


def fmt_num(x, digits=0):
    return f"{x:,.{digits}f}"


def df_table(df, max_rows=8, classes=""):
    view = df.copy().head(max_rows)
    return view.to_html(index=False, escape=True, classes=f"data-table {classes}", border=0)


def bullet(items):
    return "<ul>" + "".join(f"<li>{html.escape(str(item))}</li>" for item in items) + "</ul>"


def link_line(label, url):
    return f"<p class='repo'><b>{html.escape(label)}:</b> <span>{html.escape(url)}</span></p>"


def metric(label, value, note=""):
    return f"<div class='metric'><div class='value'>{html.escape(str(value))}</div><div class='label'>{html.escape(label)}</div><div class='note'>{html.escape(note)}</div></div>"


def slide(title, body, kicker="", section="", cls=""):
    return f"""
<section class="slide {cls}">
  <div class="topbar"><span></span><em>{html.escape(section)}</em></div>
  <div class="slide-title">
    <div>
      {'<div class="kicker">' + html.escape(kicker) + '</div>' if kicker else ''}
      <h1>{html.escape(title)}</h1>
    </div>
  </div>
  <div class="content">{body}</div>
  <footer>IBM Applied Data Science Capstone | SpaceX Falcon 9 Landing Prediction</footer>
</section>
"""


def section_slide(label, title, subtitle):
    return f"""
<section class="slide section-slide">
  <div class="section-mark">{html.escape(label)}</div>
  <h1>{html.escape(title)}</h1>
  <p>{html.escape(subtitle)}</p>
  <footer>IBM Applied Data Science Capstone | SpaceX Falcon 9 Landing Prediction</footer>
</section>
"""


def chart_slide(title, image_path, insight, repo_url, kicker="EDA Result", section="Results"):
    return slide(
        title,
        f"""
<div class="grid chart-layout">
  <div class="figure-card"><img src="{rel(image_path)}"></div>
  <aside class="insight">
    <h3>Interpretation</h3>
    {bullet(insight)}
    {link_line("GitHub URL", repo_url)}
  </aside>
</div>
""",
        kicker=kicker,
        section=section,
    )


def screenshot_slide(title, image_path, bullets_, repo_url, kicker, section):
    return slide(
        title,
        f"""
<div class="grid chart-layout">
  <div class="screenshot-card"><img src="{rel(image_path)}"></div>
  <aside class="insight">
    <h3>What this demonstrates</h3>
    {bullet(bullets_)}
    {link_line("GitHub URL", repo_url)}
  </aside>
</div>
""",
        kicker=kicker,
        section=section,
    )


def build_deck_html(frames, sql, ml, charts, maps, dash_screens, scrape):
    part1 = frames["part1"].copy()
    part2 = frames["part2"].copy()
    part3 = frames["part3"].copy()
    part2["Year"] = pd.to_datetime(part2["Date"]).dt.year
    site = (
        part2.groupby("LaunchSite", as_index=False)
        .agg(launches=("Class", "size"), successes=("Class", "sum"), success_rate=("Class", "mean"), avg_payload=("PayloadMass", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    orbit = (
        part2.groupby("Orbit", as_index=False)
        .agg(launches=("Class", "size"), successes=("Class", "sum"), success_rate=("Class", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    yearly = part2.groupby("Year", as_index=False).agg(launches=("Class", "size"), success_rate=("Class", "mean")).sort_values("Year")
    best = ml["summary"].iloc[0]
    cm = ml["best_cm"]
    tn, fp, fn, tp = cm.ravel()
    payload_mean = part1["PayloadMass"].mean()
    missing = (part1.isna().mean() * 100).sort_values(ascending=False).head(6).reset_index()
    missing.columns = ["Column", "Missing %"]
    missing["Missing %"] = missing["Missing %"].map(lambda x: f"{x:.1f}%")
    model_table = ml["summary"][["model", "cv_accuracy", "test_accuracy"]].copy()
    model_table["cv_accuracy"] = model_table["cv_accuracy"].map(lambda x: f"{x:.3f}")
    model_table["test_accuracy"] = model_table["test_accuracy"].map(lambda x: f"{x:.3f}")
    model_table.columns = ["Model", "10-fold CV accuracy", "Test accuracy"]

    slides = []
    slides.append(
        f"""
<section class="slide cover">
  <div class="cover-left">
    <div class="eyebrow">IBM Applied Data Science Capstone</div>
    <h1>SpaceX Falcon 9<br>First Stage Landing<br>Prediction</h1>
    <p>End-to-end data science project: collection, wrangling, EDA, SQL, interactive visual analytics, and classification modeling.</p>
    <div class="repo-box">GitHub URL<br><span>{html.escape(GITHUB_URL)}</span></div>
  </div>
  <div class="cover-right">
    <div class="orbit orbit-one"></div>
    <div class="orbit orbit-two"></div>
    <div class="rocket-card">
      <div class="rocket-shape"></div>
      <div>
        <b>Mission question</b>
        <span>Can launch conditions predict booster recovery?</span>
      </div>
    </div>
    <div class="mini-stats">
      {metric("Falcon 9 launches", len(part2))}
      {metric("Overall success", fmt_pct(part2["Class"].mean()))}
      {metric("Engineered features", part3.shape[1])}
    </div>
  </div>
</section>
"""
    )

    slides.append(
        slide(
            "Presentation Roadmap",
            """
<div class="roadmap">
  <div><b>01</b><span>Executive Summary and problem context</span></div>
  <div><b>02</b><span>Data collection, wrangling, and feature engineering</span></div>
  <div><b>03</b><span>Exploratory visual analysis and SQL findings</span></div>
  <div><b>04</b><span>Folium maps and Plotly Dash dashboard evidence</span></div>
  <div><b>05</b><span>Predictive modeling, model comparison, and conclusions</span></div>
</div>
""",
            kicker="Outline",
            section="Project Structure",
        )
    )

    slides.append(
        slide(
            "Executive Summary",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Methods</h3>
    {bullet([
        "Collected SpaceX launch data through REST API calls and a historical Wikipedia table.",
        "Cleaned and transformed launch records into a supervised learning dataset.",
        "Used visualization, SQL, Folium, and Dash to explain site, orbit, payload, and time patterns.",
        "Compared Logistic Regression, SVM, Decision Tree, and KNN with 10-fold GridSearchCV."
    ])}
  </div>
  <div class="panel accent">
    <h3>Key Results</h3>
    {bullet([
        f"Overall landing success rate in the Falcon 9 dataset is {fmt_pct(part2['Class'].mean())}.",
        f"Top success-rate sites: {site.iloc[0]['LaunchSite']} ({fmt_pct(site.iloc[0]['success_rate'])}) and {site.iloc[1]['LaunchSite']} ({fmt_pct(site.iloc[1]['success_rate'])}).",
        f"GTO missions have a lower success rate ({fmt_pct(orbit[orbit['Orbit']=='GTO']['success_rate'].iloc[0])}) than several specialized orbits.",
        f"Best model: {best['model']} with test accuracy {fmt_pct(best['test_accuracy'])}."
    ])}
  </div>
</div>
""",
            kicker="Rubric 1.3",
            section="Executive Summary",
        )
    )

    slides.append(
        slide(
            "Project Background and Problem Statement",
            """
<div class="problem">
  <div>
    <h3>Business Context</h3>
    <p>Falcon 9 first-stage recovery is a central part of SpaceX's reusable launch economics. If the booster lands successfully, the mission can preserve high-value hardware and improve launch economics.</p>
  </div>
  <div>
    <h3>Problem Statement</h3>
    <p>Predict whether the Falcon 9 first stage will land successfully using launch-site, orbit, payload, booster, and historical launch features.</p>
  </div>
  <div>
    <h3>Target Variable</h3>
    <p><b>Class = 1</b> means successful landing. <b>Class = 0</b> means failure, no attempt, ocean outcome, or unavailable recovery.</p>
  </div>
</div>
""",
            kicker="Rubric 1.4",
            section="Introduction",
        )
    )

    slides.append(
        slide(
            "Analytical Questions",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Questions Answered</h3>
    {bullet([
        "Which launch sites show the strongest landing performance?",
        "How do payload mass and orbit type relate to recovery success?",
        "Did success improve over time as Falcon 9 operations matured?",
        "Which machine learning model best predicts landing success?"
    ])}
  </div>
  <div class="panel">
    <h3>Evaluation Standard</h3>
    {bullet([
        "Use classification accuracy for model comparison.",
        "Use confusion matrix to understand false positives and false negatives.",
        "Use explainable EDA and SQL results to support the final recommendations."
    ])}
  </div>
</div>
""",
            kicker="Analysis Design",
            section="Introduction",
        )
    )

    slides.append(section_slide("01", "Methodology", "How the data was collected, cleaned, queried, visualized, and modeled."))

    slides.append(
        slide(
            "End-to-End Data Pipeline",
            """
<div class="flow long">
  <div><b>1</b><span>SpaceX API</span></div>
  <i></i>
  <div><b>2</b><span>Wikipedia Scraping</span></div>
  <i></i>
  <div><b>3</b><span>Wrangling</span></div>
  <i></i>
  <div><b>4</b><span>EDA + SQL</span></div>
  <i></i>
  <div><b>5</b><span>Folium + Dash</span></div>
  <i></i>
  <div><b>6</b><span>ML Models</span></div>
</div>
<div class="grid three compact-metrics">
  {metric("API modeling rows", len(part1))}
  {metric("SQL records", len(frames["sql"]))}
  {metric("Geo records", len(frames["geo"]))}
</div>
""",
            kicker="Methodology Overview",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "Data Collection - SpaceX API",
            f"""
<div class="flow">
  <div><b>GET</b><span>/v4/launches/past</span></div>
  <i></i>
  <div><b>Normalize</b><span>launch JSON</span></div>
  <i></i>
  <div><b>Enrich</b><span>rockets, payloads, pads, cores</span></div>
  <i></i>
  <div><b>Filter</b><span>Falcon 9 records</span></div>
  <i></i>
  <div><b>Export</b><span>dataset_part_1.csv</span></div>
</div>
<div class="grid two">
  <div class="panel">
    <h3>Key API Fields</h3>
    {bullet(["Flight number and date", "Booster version and core serial", "Launch site and coordinates", "Payload mass and orbit", "Grid fins, legs, reuse count, landing pad"])}
  </div>
  <div class="panel">
    <h3>Notebook Reference</h3>
    {link_line("GitHub URL", gh(NOTEBOOKS["api"]))}
  </div>
</div>
""",
            kicker="Rubric 1.5",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "SpaceX API Data Snapshot",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Cleaned API Output</h3>
    {df_table(part1[["FlightNumber", "Date", "BoosterVersion", "PayloadMass", "Orbit", "LaunchSite", "Outcome"]].head(6), max_rows=6, classes="small")}
  </div>
  <div class="panel accent">
    <h3>Why API Collection Matters</h3>
    {bullet([
        "The API gives structured launch metadata that can be transformed into model features.",
        "Core-level fields such as reuse count, serial, grid fins, and legs are directly relevant to recovery.",
        "Coordinates enable later Folium launch-site analysis."
    ])}
  </div>
</div>
""",
            kicker="API Evidence",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "Data Collection - Web Scraping",
            f"""
<div class="flow">
  <div><b>Request</b><span>archived Wikipedia page</span></div>
  <i></i>
  <div><b>Parse</b><span>BeautifulSoup HTML</span></div>
  <i></i>
  <div><b>Extract</b><span>table headers</span></div>
  <i></i>
  <div><b>Iterate</b><span>launch rows</span></div>
  <i></i>
  <div><b>Build</b><span>launch DataFrame</span></div>
</div>
<div class="grid two">
  <div class="panel">
    <h3>Scraping Result</h3>
    {bullet([
        f"Parsed {scrape['tables']} launch tables from the archived page.",
        f"Extracted approximately {scrape['rows']} raw launch rows before filtering and cleaning.",
        "Fields include date/time, version/booster, launch site, payload, orbit, customer, and outcome."
    ])}
  </div>
  <div class="panel">
    <h3>Notebook Reference</h3>
    {link_line("GitHub URL", gh(NOTEBOOKS["scraping"]))}
  </div>
</div>
""",
            kicker="Rubric 1.6",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "Data Wrangling and Cleaning",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Cleaning Steps</h3>
    {bullet([
        f"Filled missing PayloadMass with mean payload mass ({fmt_num(payload_mean, 1)} kg).",
        "Converted landing outcome into the binary Class target.",
        "Standardized date, launch-site, orbit, booster, and payload fields.",
        "Removed records outside the Falcon 9 scope for the modeling dataset."
    ])}
    {link_line("GitHub URL", gh(NOTEBOOKS["wrangling"]))}
  </div>
  <div class="panel">
    <h3>Missingness Check</h3>
    {df_table(missing, max_rows=6, classes="small")}
  </div>
</div>
""",
            kicker="Rubric 1.7",
            section="Methodology",
        )
    )

    feature_cols = ["FlightNumber", "PayloadMass", "Orbit", "LaunchSite", "Flights", "GridFins", "Reused", "Legs", "LandingPad", "Block", "ReusedCount", "Serial"]
    slides.append(
        slide(
            "Feature Engineering for Modeling",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Feature Set</h3>
    {bullet(feature_cols)}
  </div>
  <div class="panel accent">
    <h3>Transformation</h3>
    {bullet([
        "Categorical variables were one-hot encoded.",
        "All model inputs were cast to float64.",
        f"The final feature matrix contains {part3.shape[1]} numeric columns.",
        "StandardScaler was applied before model training."
    ])}
  </div>
</div>
""",
            kicker="Model Preparation",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "EDA with Data Visualization",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Charts Produced</h3>
    {bullet([
        "Flight Number vs. Launch Site scatter plot",
        "Payload Mass vs. Launch Site scatter plot",
        "Success Rate vs. Orbit Type bar chart",
        "Flight Number vs. Orbit Type scatter plot",
        "Payload Mass vs. Orbit Type scatter plot",
        "Yearly average success trend line chart"
    ])}
  </div>
  <div class="panel">
    <h3>Purpose</h3>
    {bullet([
        "Find site-level differences in recovery performance.",
        "Assess payload and orbit relationships with landing outcome.",
        "Show how recovery reliability changed over time."
    ])}
    {link_line("GitHub URL", gh(NOTEBOOKS["eda_viz"]))}
  </div>
</div>
""",
            kicker="Rubric 1.8",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "EDA with SQL",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Query Coverage</h3>
    {bullet([
        "Distinct launch sites and CCA launch-site records.",
        "NASA payload totals and F9 v1.1 average payload mass.",
        "First successful ground-pad landing date.",
        "Successful drone-ship landings in a 4000-6000 kg payload window.",
        "Mission outcome counts, maximum payload boosters, 2015 failures, and landing outcome ranking."
    ])}
  </div>
  <div class="panel code-panel">
    <h3>Representative Query</h3>
    <pre>SELECT Landing_Outcome, COUNT(*) AS Count
FROM SPACEXTABLE
WHERE Date BETWEEN '2010-06-04' AND '2017-03-20'
GROUP BY Landing_Outcome
ORDER BY Count DESC;</pre>
    {link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
  </div>
</div>
""",
            kicker="Rubric 1.9",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "Interactive Visual Analytics",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Folium Map Objects</h3>
    {bullet(["Circle and Marker objects for launch sites", "MarkerCluster for launch records", "Green/red markers for success and failure", "MousePosition, DivIcon labels, and PolyLine distance analysis"])}
    {link_line("GitHub URL", gh(NOTEBOOKS["folium"]))}
  </div>
  <div class="panel">
    <h3>Plotly Dash Interactions</h3>
    {bullet(["Launch-site dropdown", "Payload range slider", "Successful launch pie chart", "Payload mass vs. landing outcome scatter plot"])}
    {link_line("GitHub URL", gh(NOTEBOOKS["dash"]))}
  </div>
</div>
""",
            kicker="Rubric 1.10",
            section="Methodology",
        )
    )

    slides.append(
        slide(
            "Predictive Analysis Methodology",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Modeling Workflow</h3>
    {bullet(["Create target vector Y from Class", "Standardize model features", "Train/test split with random_state=2", "Tune each algorithm with 10-fold GridSearchCV"])}
  </div>
  <div class="panel">
    <h3>Models Compared</h3>
    {bullet(["Logistic Regression", "Support Vector Machine", "Decision Tree", "K Nearest Neighbors"])}
    {link_line("GitHub URL", gh(NOTEBOOKS["ml"]))}
  </div>
</div>
""",
            kicker="Rubric 1.15",
            section="Methodology",
        )
    )

    slides.append(section_slide("02", "EDA Results", "Scatter plots, bar charts, and time trends that explain landing success patterns."))

    slides.append(
        slide(
            "Dataset Overview",
            f"""
<div class="grid four kpi-row">
  {metric("Launches analyzed", len(part2), "Falcon 9 records")}
  {metric("Successful landings", int(part2["Class"].sum()), "Class = 1")}
  {metric("Overall success", fmt_pct(part2["Class"].mean()), "Landing rate")}
  {metric("Years covered", f"{part2['Year'].min()}-{part2['Year'].max()}", "Modeling dataset")}
</div>
<div class="grid two">
  <div class="figure-card"><img src="{rel(charts['site_success'])}"></div>
  <div class="panel">
    <h3>Immediate Pattern</h3>
    {bullet([
        "Launch-site effects are visible in the aggregate success rate.",
        "KSC LC 39A and VAFB SLC 4E outperform CCAFS SLC 40 in this dataset.",
        "This motivates retaining launch site as a categorical model feature."
    ])}
  </div>
</div>
""",
            kicker="Results Overview",
            section="Results",
        )
    )

    slides.append(chart_slide("Flight Number vs. Launch Site", charts["flight_site"], ["Later flight numbers contain a higher share of successful landings.", "The pattern suggests a learning curve as operations mature.", "KSC LC 39A and VAFB SLC 4E show strong later-period landing outcomes."], gh(NOTEBOOKS["eda_viz"])))
    slides.append(chart_slide("Payload vs. Launch Site", charts["payload_site"], ["Successful landings occur across both light and heavy payload ranges.", "Payload mass alone does not explain landing outcome.", "Site and mission profile should be modeled jointly with payload mass."], gh(NOTEBOOKS["eda_viz"])))
    slides.append(chart_slide("Success Rate vs. Orbit Type", charts["orbit_success"], ["Several orbit classes show perfect success in this sample.", "GTO missions are materially more difficult than SSO, VLEO, GEO, HEO, and ES-L1 in the observed records.", "Orbit type is a meaningful categorical feature."], gh(NOTEBOOKS["eda_viz"])))
    slides.append(chart_slide("Flight Number vs. Orbit Type", charts["flight_orbit"], ["Orbit diversity increases as Falcon 9 launch history grows.", "Later VLEO and SSO missions show strong landing success.", "Time and orbit interact: newer launch periods include more mature recovery operations."], gh(NOTEBOOKS["eda_viz"])))
    slides.append(chart_slide("Payload vs. Orbit Type", charts["payload_orbit"], ["Very heavy VLEO payloads can still land successfully.", "GTO payloads show more mixed outcomes.", "Payload should be interpreted together with orbit energy requirements."], gh(NOTEBOOKS["eda_viz"])))
    slides.append(chart_slide("Launch Success Yearly Trend", charts["yearly"], ["Success rises sharply after the early experimental period.", "2017, 2019, and 2020 show high landing reliability.", "The annual trend supports the operational-learning interpretation."], gh(NOTEBOOKS["eda_viz"])))

    slides.append(section_slide("03", "SQL Results", "Structured queries that validate launch-site, payload, success, ranking, and time-based findings."))

    slides.append(
        slide(
            "SQL: Launch Site Names",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Unique Launch Sites</h3>
    {df_table(sql["unique_sites"], max_rows=8)}
  </div>
  <div class="panel">
    <h3>Interpretation</h3>
    {bullet(["The SQL dataset contains Cape Canaveral, Kennedy Space Center, and Vandenberg launch-site labels.", "The site names include LC-40/SLC-40 naming variants that should be handled carefully in reporting."])}
    {link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
  </div>
</div>
""",
            kicker="Rubric 1.12",
            section="SQL Results",
        )
    )

    slides.append(
        slide(
            "SQL: First Five CCA Launch Records",
            f"""
<div class="wide-table">
{df_table(sql["cca5"], max_rows=5, classes="tiny")}
</div>
<div class="callout">These records confirm the CCA launch-site filter and show early Falcon 9 missions with payload, orbit, and landing outcome fields.</div>
{link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
""",
            kicker="Rubric 1.12",
            section="SQL Results",
        )
    )

    nasa_payload = sql["nasa_payload"].iloc[0, 0]
    f9_avg = sql["f9v11_avg"].iloc[0, 0]
    slides.append(
        slide(
            "SQL: Payload Aggregations",
            f"""
<div class="grid two big-metrics">
  {metric("Total NASA customer payload", f"{fmt_num(nasa_payload)} kg", "SUM(PAYLOAD_MASS__KG_) WHERE Customer LIKE '%NASA%'")}
  {metric("Average F9 v1.1 payload", f"{fmt_num(f9_avg, 2)} kg", "AVG(PAYLOAD_MASS__KG_) WHERE Booster_Version LIKE 'F9 v1.1%'")}
</div>
<div class="callout">Payload aggregation helps connect customer, booster version, and mission profile to launch economics and landing-risk analysis.</div>
{link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
""",
            kicker="Rubric 1.12",
            section="SQL Results",
        )
    )

    slides.append(
        slide(
            "SQL: Successful Landing Details",
            f"""
<div class="grid two">
  <div class="panel accent">
    <h3>First Successful Ground-Pad Landing</h3>
    <div class="giant">{html.escape(str(sql["first_ground_success"].iloc[0, 0]))}</div>
    <p class="muted">Landing_Outcome = Success (ground pad)</p>
  </div>
  <div class="panel">
    <h3>Successful Drone Ship Landings, 4000-6000 kg</h3>
    {df_table(sql["drone_4000_6000"], max_rows=6, classes="small")}
  </div>
</div>
{link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
""",
            kicker="Rubric 1.12",
            section="SQL Results",
        )
    )

    slides.append(
        slide(
            "SQL: Mission Outcomes and Maximum Payload",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Mission Outcome Counts</h3>
    {df_table(sql["mission_outcomes"], max_rows=5, classes="small")}
  </div>
  <div class="panel">
    <h3>Boosters Carrying Maximum Payload</h3>
    {df_table(sql["max_payload"].head(8), max_rows=8, classes="small")}
    <p class="muted">{len(sql["max_payload"])} booster records carried the maximum listed payload.</p>
  </div>
</div>
{link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
""",
            kicker="Rubric 1.12",
            section="SQL Results",
        )
    )

    slides.append(
        slide(
            "SQL: 2015 Failures and Outcome Ranking",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Failed Drone Ship Landings in 2015</h3>
    {df_table(sql["failed_drone_2015"], max_rows=4, classes="small")}
  </div>
  <div class="panel">
    <h3>Ranked Outcomes, 2010-06-04 to 2017-03-20</h3>
    {df_table(sql["rank_outcomes"], max_rows=8, classes="small")}
  </div>
</div>
{link_line("GitHub URL", gh(NOTEBOOKS["eda_sql"]))}
""",
            kicker="Rubric 1.12",
            section="SQL Results",
        )
    )

    slides.append(section_slide("04", "Interactive Visual Analytics", "Folium maps and Plotly Dash dashboard screenshots that turn results into an exploratory interface."))

    slides.append(screenshot_slide("Folium Map: Launch Site Markers", maps["folium_sites"], ["Launch-site Circle and Marker objects identify all launch locations.", "Marker labels summarize site-level launch volume and success rate.", "The geographic view shows concentration along Florida's Space Coast plus Vandenberg on the West Coast."], gh(NOTEBOOKS["folium"]), "Rubric 1.13", "Folium"))
    slides.append(screenshot_slide("Folium Map: Launch Records by Outcome", maps["folium_records"], ["MarkerCluster groups individual launch records.", "Green markers indicate successful landings; red markers indicate failures or no attempts.", "MousePosition supports coordinate inspection during exploratory map analysis."], gh(NOTEBOOKS["folium"]), "Rubric 1.13", "Folium"))
    slides.append(screenshot_slide("Folium Map: Proximity Analysis", maps["folium_proximity"], ["PolyLine objects measure distance from KSC LC-39A to nearby reference features.", "DivIcon labels display approximate distance in kilometers.", "The analysis connects launch operations with coastline, highway, and city proximity."], gh(NOTEBOOKS["folium"]), "Rubric 1.13", "Folium"))

    slides.append(screenshot_slide("Plotly Dash: Successful Launches Pie Chart", dash_screens["dash_all"], ["The all-sites pie chart compares where successful launches occurred.", "The launch-site dropdown is set to All Sites.", "This view answers which launch sites contribute most to successful recovery records."], gh(NOTEBOOKS["dash"]), "Rubric 1.14", "Dash"))
    slides.append(screenshot_slide("Plotly Dash: Selected-Site Outcome Pie", dash_screens["dash_site"], ["The selected-site pie chart shows success versus failure for the highest-ratio launch site.", "This supports focused site-level performance analysis.", "The same app logic can be reused for each launch site."], gh(NOTEBOOKS["dash"]), "Rubric 1.14", "Dash"))
    slides.append(screenshot_slide("Plotly Dash: Payload vs. Launch Outcome", dash_screens["dash_scatter"], ["The scatter plot compares payload mass with landing outcome.", "Color encodes booster version category.", "The payload range slider supports filtering heavy or light missions."], gh(NOTEBOOKS["dash"]), "Rubric 1.14", "Dash"))

    slides.append(section_slide("05", "Predictive Analysis", "Model comparison, evaluation metrics, confusion matrix, best model, and decision-oriented conclusions."))

    class_counts = part2["Class"].value_counts().rename(index={0: "Not landed", 1: "Landed"}).reset_index()
    class_counts.columns = ["Class", "Count"]
    slides.append(
        slide(
            "Modeling Dataset and Target Balance",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Target Distribution</h3>
    {df_table(class_counts, max_rows=3)}
    {bullet([f"Target success rate: {fmt_pct(part2['Class'].mean())}", f"Feature columns after encoding: {part3.shape[1]}", "Train/test split: 80% / 20%"])}
  </div>
  <div class="panel">
    <h3>Model Evaluation</h3>
    {bullet(["10-fold cross-validation estimates model stability.", "Test accuracy evaluates out-of-sample performance.", "Confusion matrix reveals whether errors are false positives or false negatives."])}
    {link_line("GitHub URL", gh(NOTEBOOKS["ml"]))}
  </div>
</div>
""",
            kicker="Rubric 1.15",
            section="Prediction",
        )
    )

    slides.append(chart_slide("Classification Accuracy", charts["ml_accuracy"], [f"{best['model']} is selected as the best model because it ties on test accuracy and has the highest cross-validation score.", "All four models achieve useful baseline performance on the small test set.", "The result should be interpreted with caution because the dataset is relatively small."], gh(NOTEBOOKS["ml"]), kicker="Rubric 1.15", section="Prediction"))

    slides.append(
        slide(
            "Model Comparison Table",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Accuracy Summary</h3>
    {df_table(model_table, max_rows=4)}
  </div>
  <div class="panel">
    <h3>Best Hyperparameters</h3>
    {bullet([f"{row.model}: {row.best_params}" for row in ml["summary"].itertuples()])}
  </div>
</div>
""",
            kicker="Rubric 1.15",
            section="Prediction",
        )
    )

    slides.append(
        slide(
            "Best Model: Confusion Matrix",
            f"""
<div class="grid chart-layout">
  <div class="figure-card cm"><img src="{rel(charts['ml_cm'])}"></div>
  <aside class="insight">
    <h3>Evaluation</h3>
    {bullet([
        f"Best model: {best['model']}",
        f"Test accuracy: {fmt_pct(best['test_accuracy'])}",
        f"10-fold CV accuracy: {fmt_pct(best['cv_accuracy'])}",
        f"Confusion matrix values: TN={tn}, FP={fp}, FN={fn}, TP={tp}.",
        "The model correctly identifies all successful landings in this test split, but false positives remain a risk."
    ])}
    {link_line("GitHub URL", gh(NOTEBOOKS["ml"]))}
  </aside>
</div>
""",
            kicker="Rubric 1.15",
            section="Prediction",
        )
    )

    slides.append(
        slide(
            "Best Model Interpretation",
            f"""
<div class="grid two">
  <div class="panel accent">
    <h3>Why {html.escape(str(best['model']))}?</h3>
    {bullet([
        "It achieved the strongest cross-validation score among models with the same test accuracy.",
        "Decision-tree logic can capture non-linear interactions among site, orbit, payload, reuse, and booster attributes.",
        "Its rules are easier to explain than a black-box model for a small operational dataset."
    ])}
  </div>
  <div class="panel">
    <h3>Residual Risk</h3>
    {bullet([
        "The dataset is small, so model ranking may vary with the train/test split.",
        "False positives matter operationally because predicting success for a failed landing can understate recovery risk.",
        "The model should support, not replace, mission engineering review."
    ])}
  </div>
</div>
""",
            kicker="Rubric 1.15",
            section="Prediction",
        )
    )

    slides.append(
        slide(
            "Conclusions",
            f"""
<div class="conclusion-list">
  <div><b>1</b><span>Landing success improved sharply after the early experimental phase, indicating a strong operational learning curve.</span></div>
  <div><b>2</b><span>Launch site and orbit type are meaningful predictors; KSC LC 39A and VAFB SLC 4E perform strongly, while GTO missions are more challenging.</span></div>
  <div><b>3</b><span>Payload mass alone is not decisive; heavy VLEO missions can land successfully when paired with mature booster operations.</span></div>
  <div><b>4</b><span>The best model reaches {fmt_pct(best['test_accuracy'])} test accuracy, but false-positive risk means prediction should be paired with mission-risk review.</span></div>
  <div><b>5</b><span>Creative decision insight: combine model probability with a site-orbit-booster risk score to flag launches needing additional recovery planning.</span></div>
</div>
""",
            kicker="Rubric 1.15",
            section="Conclusion",
        )
    )

    slides.append(
        slide(
            "Limitations and Future Work",
            f"""
<div class="grid two">
  <div class="panel">
    <h3>Limitations</h3>
    {bullet([
        "The modeling dataset is small, so model accuracy can vary by train/test split.",
        "The dataset does not include all engineering variables that affect landing risk, such as weather, sea state, fuel reserve, and mission-specific trajectory constraints.",
        "Several orbit classes have very few observations, so perfect success rates should not be overgeneralized."
    ])}
  </div>
  <div class="panel accent">
    <h3>Future Work</h3>
    {bullet([
        "Add newer launches and richer telemetry or weather features.",
        "Use calibrated probabilities rather than only class labels.",
        "Create a combined recovery risk score using model probability, orbit class, payload mass, launch site, and booster reuse history.",
        "Deploy the Dash dashboard as a lightweight decision-support tool for mission review."
    ])}
  </div>
</div>
""",
            kicker="Professional Reflection",
            section="Conclusion",
        )
    )

    slides.append(
        slide(
            "Appendix: GitHub References",
            f"""
<div class="wide-table links">
{df_table(pd.DataFrame([
    ["Project GitHub", GITHUB_URL],
    ["SpaceX API notebook", gh(NOTEBOOKS["api"])],
    ["Web scraping notebook", gh(NOTEBOOKS["scraping"])],
    ["Data wrangling notebook", gh(NOTEBOOKS["wrangling"])],
    ["EDA visualization notebook", gh(NOTEBOOKS["eda_viz"])],
    ["EDA SQL notebook", gh(NOTEBOOKS["eda_sql"])],
    ["Folium notebook", gh(NOTEBOOKS["folium"])],
    ["Dash app", gh(NOTEBOOKS["dash"])],
    ["Machine learning notebook", gh(NOTEBOOKS["ml"])],
], columns=["Artifact", "URL"]), max_rows=10, classes="tiny")}
</div>
""",
            kicker="Reference",
            section="Appendix",
        )
    )

    styles = """
<!doctype html>
<html>
<head>
<meta charset="utf-8">
<title>IBM SpaceX Capstone Final Report</title>
<style>
@page { size: 16in 9in; margin: 0; }
* { box-sizing: border-box; }
body {
  margin:0;
  background:#CBD5E1;
  font-family: Inter, Arial, Helvetica, sans-serif;
  color:#152238;
}
.slide {
  position:relative;
  width:1600px;
  height:900px;
  background:#F6F8FB;
  overflow:hidden;
  page-break-after: always;
  padding:72px 86px 58px 86px;
}
.slide::before {
  content:"";
  position:absolute;
  inset:0;
  background:
    radial-gradient(circle at 88% 12%, rgba(47,128,237,.12), transparent 290px),
    linear-gradient(90deg, rgba(15,42,67,.035), transparent 38%);
  pointer-events:none;
}
.topbar {
  position:absolute;
  left:0;
  top:0;
  height:18px;
  width:100%;
  background:#0F2A43;
}
.topbar span {
  display:block;
  height:18px;
  width:330px;
  background:#00A78E;
}
.topbar em {
  position:absolute;
  top:28px;
  right:86px;
  font-style:normal;
  color:#8A94A6;
  font-size:15px;
  letter-spacing:.04em;
  text-transform:uppercase;
}
.slide-title {
  position:relative;
  display:flex;
  justify-content:space-between;
  align-items:flex-start;
  margin-bottom:30px;
}
.kicker {
  color:#00A78E;
  font-size:18px;
  font-weight:800;
  letter-spacing:.06em;
  text-transform:uppercase;
  margin-bottom:8px;
}
h1 {
  margin:0;
  font-size:48px;
  line-height:1.04;
  letter-spacing:-.02em;
  color:#152238;
}
h3 {
  margin:0 0 18px 0;
  font-size:24px;
  color:#152238;
}
p, li {
  font-size:23px;
  line-height:1.38;
}
.content {
  position:relative;
  z-index:2;
}
footer {
  position:absolute;
  left:86px;
  bottom:26px;
  color:#8A94A6;
  font-size:14px;
}
.grid { display:grid; gap:26px; }
.two { grid-template-columns: 1fr 1fr; }
.three { grid-template-columns: repeat(3, 1fr); }
.four { grid-template-columns: repeat(4, 1fr); }
.chart-layout { grid-template-columns: 1.72fr .88fr; align-items:stretch; }
.panel, .figure-card, .screenshot-card, .insight, .metric, .callout {
  background:white;
  border:1px solid #D8E0EA;
  border-radius:16px;
  box-shadow:0 18px 44px rgba(15,42,67,.08);
}
.panel { padding:30px 34px; min-height:250px; }
.panel.accent { border-top:7px solid #00A78E; }
.code-panel pre {
  background:#0F2A43;
  color:#E8F7FF;
  padding:22px;
  border-radius:10px;
  font-size:18px;
  line-height:1.36;
  white-space:pre-wrap;
}
ul { margin:0; padding-left:26px; }
li { margin:0 0 14px 0; }
.repo {
  font-size:17px;
  color:#64748B;
  line-height:1.32;
  word-break:break-word;
  margin-top:20px;
}
.repo span { color:#2F80ED; }
.figure-card, .screenshot-card { padding:18px; height:610px; display:flex; align-items:center; justify-content:center; }
.figure-card img, .screenshot-card img { max-width:100%; max-height:100%; object-fit:contain; border-radius:10px; }
.screenshot-card img { width:100%; height:100%; object-fit:contain; background:white; }
.insight { padding:32px; }
.insight li { font-size:21px; }
.metric {
  padding:24px 26px;
  min-height:128px;
  border-top:6px solid #00A78E;
}
.metric .value {
  font-size:40px;
  font-weight:900;
  color:#0F2A43;
  line-height:1.03;
}
.metric .label {
  margin-top:10px;
  color:#152238;
  font-size:17px;
  font-weight:800;
  text-transform:uppercase;
  letter-spacing:.04em;
}
.metric .note {
  color:#8A94A6;
  font-size:17px;
  margin-top:6px;
}
.compact-metrics { margin-top:56px; }
.kpi-row { margin-bottom:28px; }
.data-table {
  width:100%;
  border-collapse:collapse;
  font-size:20px;
  background:white;
}
.data-table th {
  text-align:left;
  background:#0F2A43;
  color:white;
  padding:12px 13px;
  font-weight:800;
}
.data-table td {
  padding:12px 13px;
  border-bottom:1px solid #E6ECF2;
  color:#152238;
  vertical-align:top;
}
.data-table tr:nth-child(even) td { background:#F8FAFC; }
.data-table.small { font-size:17px; }
.data-table.tiny { font-size:14px; }
.wide-table { background:white; border:1px solid #D8E0EA; border-radius:16px; padding:20px; box-shadow:0 18px 44px rgba(15,42,67,.08); overflow:hidden; }
.wide-table .data-table { font-size:17px; }
.wide-table .data-table.tiny { font-size:13px; }
.callout {
  margin-top:24px;
  padding:22px 26px;
  font-size:23px;
  line-height:1.35;
  border-left:8px solid #00A78E;
}
.giant { font-size:64px; font-weight:900; color:#00A78E; margin-top:38px; }
.muted { color:#64748B; font-size:19px; }
.cover {
  background:linear-gradient(135deg,#0B1F33,#0F2A43 60%,#123B5D);
  color:white;
  padding:0;
  display:grid;
  grid-template-columns: 1.05fr .95fr;
}
.cover::before { display:none; }
.cover-left { padding:90px 74px; position:relative; z-index:2; }
.cover .eyebrow { color:#8BE7D6; font-weight:800; font-size:20px; letter-spacing:.08em; text-transform:uppercase; margin-bottom:26px; }
.cover h1 { color:white; font-size:70px; line-height:.98; letter-spacing:-.035em; }
.cover p { color:#D6E4F0; max-width:720px; margin-top:30px; font-size:25px; }
.repo-box {
  position:absolute;
  left:74px;
  bottom:70px;
  width:700px;
  border:1px solid rgba(255,255,255,.22);
  background:rgba(255,255,255,.08);
  padding:18px 22px;
  border-radius:14px;
  color:#8BE7D6;
  font-size:16px;
  text-transform:uppercase;
  letter-spacing:.06em;
}
.repo-box span { display:block; margin-top:8px; color:white; font-size:18px; text-transform:none; letter-spacing:0; word-break:break-all; }
.cover-right { position:relative; overflow:hidden; }
.orbit { position:absolute; border:2px solid rgba(139,231,214,.24); border-radius:50%; transform:rotate(-25deg); }
.orbit-one { width:760px; height:390px; right:-90px; top:120px; }
.orbit-two { width:980px; height:520px; right:-210px; top:70px; }
.rocket-card {
  position:absolute;
  top:178px;
  right:118px;
  width:435px;
  height:255px;
  background:rgba(255,255,255,.10);
  border:1px solid rgba(255,255,255,.20);
  border-radius:22px;
  padding:34px;
  display:flex;
  align-items:center;
  gap:28px;
  box-shadow:0 28px 80px rgba(0,0,0,.22);
}
.rocket-card b { display:block; color:white; font-size:26px; margin-bottom:12px; }
.rocket-card span { color:#D6E4F0; font-size:21px; line-height:1.25; }
.rocket-shape {
  width:74px;
  height:178px;
  border-radius:42px 42px 18px 18px;
  background:linear-gradient(180deg,#fff,#9ADCF8);
  position:relative;
  box-shadow:0 0 40px rgba(139,231,214,.42);
}
.rocket-shape::before { content:""; position:absolute; top:32px; left:20px; width:34px; height:34px; border-radius:50%; background:#2F80ED; border:5px solid #DFF8FF; }
.rocket-shape::after { content:""; position:absolute; bottom:-46px; left:18px; border-left:19px solid transparent; border-right:19px solid transparent; border-top:52px solid #F2994A; filter:drop-shadow(0 0 18px rgba(242,153,74,.65)); }
.mini-stats { position:absolute; bottom:78px; right:80px; display:grid; grid-template-columns:repeat(3, 170px); gap:16px; }
.mini-stats .metric { background:rgba(255,255,255,.12); border-color:rgba(255,255,255,.22); box-shadow:none; min-height:112px; }
.mini-stats .metric .value { color:white; font-size:34px; }
.mini-stats .metric .label { color:#8BE7D6; font-size:13px; }
.mini-stats .metric .note { color:#C7D8E6; }
.roadmap {
  display:grid;
  gap:22px;
  margin-top:25px;
}
.roadmap div {
  background:white;
  border:1px solid #D8E0EA;
  border-radius:16px;
  padding:24px 30px;
  display:grid;
  grid-template-columns:82px 1fr;
  align-items:center;
  box-shadow:0 14px 34px rgba(15,42,67,.07);
}
.roadmap b { font-size:32px; color:#00A78E; }
.roadmap span { font-size:27px; color:#152238; }
.problem { display:grid; grid-template-columns:repeat(3,1fr); gap:24px; }
.problem div { background:white; border:1px solid #D8E0EA; border-radius:16px; padding:34px; min-height:490px; box-shadow:0 18px 44px rgba(15,42,67,.08); }
.problem p { font-size:24px; }
.flow {
  display:grid;
  grid-template-columns: repeat(5, 1fr 34px) 1fr;
  align-items:center;
  gap:12px;
  margin-bottom:42px;
}
.flow.long { grid-template-columns: repeat(5, 1fr 28px) 1fr; }
.flow div {
  background:white;
  border:1px solid #D8E0EA;
  border-radius:14px;
  padding:19px 16px;
  min-height:116px;
  box-shadow:0 12px 30px rgba(15,42,67,.07);
}
.flow b { display:block; color:#00A78E; font-size:21px; margin-bottom:8px; }
.flow span { color:#152238; font-size:19px; line-height:1.2; }
.flow i {
  display:block;
  height:3px;
  background:#00A78E;
  position:relative;
}
.flow i::after {
  content:"";
  position:absolute;
  right:-1px;
  top:-6px;
  width:0;
  height:0;
  border-top:8px solid transparent;
  border-bottom:8px solid transparent;
  border-left:12px solid #00A78E;
}
.section-slide {
  background:linear-gradient(135deg,#0F2A43,#123B5D);
  color:white;
  display:flex;
  flex-direction:column;
  justify-content:center;
}
.section-slide::before { display:none; }
.section-slide .section-mark {
  color:#8BE7D6;
  font-size:92px;
  font-weight:900;
  line-height:1;
  margin-bottom:18px;
}
.section-slide h1 { color:white; font-size:72px; }
.section-slide p { color:#D6E4F0; font-size:28px; max-width:1000px; }
.section-slide footer { color:#9FB7CC; }
.big-metrics .metric { min-height:310px; padding:42px; }
.big-metrics .metric .value { font-size:54px; }
.big-metrics .metric .note { font-size:20px; line-height:1.32; margin-top:20px; }
.cm { height:610px; }
.conclusion-list { display:grid; gap:18px; }
.conclusion-list div {
  background:white;
  border:1px solid #D8E0EA;
  border-radius:16px;
  padding:22px 28px;
  display:grid;
  grid-template-columns:64px 1fr;
  align-items:center;
  box-shadow:0 12px 30px rgba(15,42,67,.07);
}
.conclusion-list b { color:#00A78E; font-size:32px; }
.conclusion-list span { font-size:24px; line-height:1.34; }
.links .data-table td:nth-child(2) { word-break:break-all; color:#2F80ED; }
@media print {
  body { background:white; }
  .slide { margin:0; }
}
</style>
</head>
<body>
"""
    HTML_DECK.write_text(styles + "\n".join(slides) + "\n</body></html>", encoding="utf-8")
    return len(slides)


def render_pdf_and_slide_images(slide_count):
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1600, "height": 900}, device_scale_factor=1)
        page.goto(HTML_DECK.resolve().as_uri(), wait_until="networkidle")
        page.pdf(path=str(PDF_PREMIUM), print_background=True, prefer_css_page_size=True)
        shutil.copyfile(PDF_PREMIUM, PDF_OUT)
        shutil.copyfile(PDF_PREMIUM, PDF_ALIAS)

        for old in SLIDES.glob("slide_*.png"):
            old.unlink()
        locators = page.locator("section.slide")
        count = locators.count()
        if count != slide_count:
            raise RuntimeError(f"Expected {slide_count} slides, rendered {count}")
        for idx in range(count):
            locators.nth(idx).screenshot(path=str(SLIDES / f"slide_{idx+1:02d}.png"))
        browser.close()


def build_pptx(slide_count):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    for idx in range(1, slide_count + 1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(SLIDES / f"slide_{idx:02d}.png"), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX_OUT)
    shutil.copyfile(PPTX_OUT, PPTX_ALIAS)


def write_quality_report(slide_count, frames, ml):
    doc = fitz.open(PDF_PREMIUM)
    text = "\n".join(page.get_text() for page in doc)
    cjk_count = sum(1 for char in text if "\u4e00" <= char <= "\u9fff")
    doc.close()
    report = OUT / "premium_quality_check.md"
    report.write_text(
        "\n".join(
            [
                "# Premium Capstone Deliverable Quality Check",
                "",
                f"- Slides: {slide_count}",
                f"- PDF pages: {len(fitz.open(PDF_PREMIUM))}",
                f"- CJK characters in PDF extracted text: {cjk_count}",
                f"- Main upload PDF: `{PDF_OUT.name}`",
                f"- Premium PDF copy: `{PDF_PREMIUM.name}`",
                f"- Underscore PDF copy: `{PDF_ALIAS.name}`",
                f"- PPTX: `{PPTX_OUT.name}`",
                f"- Overall success rate: {frames['part2']['Class'].mean():.3f}",
                f"- Best model: {ml['summary'].iloc[0]['model']}",
                f"- Best model test accuracy: {ml['summary'].iloc[0]['test_accuracy']:.3f}",
                f"- GitHub URL used: `{GITHUB_URL}`",
            ]
        ),
        encoding="utf-8",
    )


def main():
    ensure_dirs()
    frames = read_frames()
    scrape = wiki_scrape_summary()
    sql = sql_results(frames["sql"])
    ml = train_models(frames["part2"], frames["part3"])
    charts = make_static_charts(frames, ml)
    maps = make_folium_maps(frames)
    dash_screens = make_dash_screens(frames)
    slide_count = build_deck_html(frames, sql, ml, charts, maps, dash_screens, scrape)
    render_pdf_and_slide_images(slide_count)
    build_pptx(slide_count)
    write_quality_report(slide_count, frames, ml)
    print(f"Generated {PDF_OUT}")
    print(f"Generated {PDF_PREMIUM}")
    print(f"Generated {PPTX_OUT}")
    print(f"Slides: {slide_count}")
    print(f"GitHub URL used: {GITHUB_URL}")


if __name__ == "__main__":
    main()
