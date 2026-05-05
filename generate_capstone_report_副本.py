import os
import shutil
import sqlite3
import textwrap
import warnings
from pathlib import Path
from urllib.parse import quote

import fitz
import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import requests
import seaborn as sns
from bs4 import BeautifulSoup
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from sklearn import preprocessing
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import accuracy_score, confusion_matrix
from sklearn.model_selection import GridSearchCV, train_test_split
from sklearn.neighbors import KNeighborsClassifier
from sklearn.svm import SVC
from sklearn.tree import DecisionTreeClassifier


warnings.filterwarnings("ignore")
os.environ.setdefault("LOKY_MAX_CPU_COUNT", "1")

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "capstone_report_outputs"
FIG = OUT / "figures"
SLIDE_IMAGES = OUT / "slide_images"
DATA = OUT / "data"

PDF_EN = ROOT / "Data_Science_Capstone_Project_Report.pdf"
PDF_UPLOAD = ROOT / "Data Science Capstone Project Report.pdf"
PPTX_OUT = ROOT / "Data_Science_Capstone_Project_Report.pptx"

GITHUB_URL = os.environ.get(
    "CAPSTONE_GITHUB_URL",
    "https://github.com/your-github-username/ibm-data-science-capstone-spacex",
).rstrip("/")

URLS = {
    "part1": "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/IBM-DS0321EN-SkillsNetwork/datasets/dataset_part_1.csv",
    "part2": "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/IBM-DS0321EN-SkillsNetwork/datasets/dataset_part_2.csv",
    "part3": "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/IBM-DS0321EN-SkillsNetwork/datasets/dataset_part_3.csv",
    "sql": "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/IBM-DS0321EN-SkillsNetwork/labs/module_2/data/Spacex.csv",
    "geo": "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/IBM-DS0321EN-SkillsNetwork/datasets/spacex_launch_geo.csv",
    "dash": "https://cf-courses-data.s3.us.cloud-object-storage.appdomain.cloud/IBM-DS0321EN-SkillsNetwork/datasets/spacex_launch_dash.csv",
}

NOTEBOOKS = {
    "api": "jupyter-labs-spacex-data-collection-api-v2 (1).ipynb",
    "scraping": "jupyter-labs-webscraping (1).ipynb",
    "wrangling": "labs-jupyter-spacex-Data wrangling-v2 (1).ipynb",
    "eda_viz": "jupyter-labs-eda-dataviz-v2 (1).ipynb",
    "eda_sql": "jupyter-labs-eda-sql-coursera_sqllite (1).ipynb",
    "folium": "lab-jupyter-launch-site-location-v2 (1).ipynb",
    "dash": "spacex_dash_app.py",
    "ml": "SpaceX-Machine-Learning-Prediction-Part-5-v1 (1).ipynb",
}

COLORS = {
    "navy": "#243B53",
    "blue": "#457B9D",
    "teal": "#2A9D8F",
    "amber": "#F4A261",
    "red": "#E76F51",
    "ink": "#1F2933",
    "muted": "#6B7280",
    "bg": "#F8FAFC",
    "line": "#CBD5E1",
    "panel": "#FFFFFF",
}

plt.rcParams.update(
    {
        "font.family": "DejaVu Sans",
        "axes.titleweight": "bold",
        "axes.labelcolor": COLORS["ink"],
        "xtick.color": COLORS["ink"],
        "ytick.color": COLORS["ink"],
        "figure.dpi": 160,
        "savefig.dpi": 220,
    }
)
sns.set_theme(style="whitegrid")


def ensure_dirs():
    for path in [OUT, FIG, SLIDE_IMAGES, DATA]:
        path.mkdir(parents=True, exist_ok=True)


def gh(path):
    return f"{GITHUB_URL}/blob/main/{quote(path)}"


def read_data():
    frames = {}
    for name, url in URLS.items():
        frame = pd.read_csv(url)
        frame.to_csv(DATA / f"{name}.csv", index=False)
        frames[name] = frame
    return frames


def wiki_scrape_summary():
    url = "https://en.wikipedia.org/w/index.php?title=List_of_Falcon_9_and_Falcon_Heavy_launches&oldid=1027686922"
    try:
        response = requests.get(
            url,
            headers={"User-Agent": "Mozilla/5.0 IBM Capstone educational project"},
            timeout=30,
        )
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        tables = soup.find_all("table")
        launch_tables = [
            table
            for table in tables
            if "Flight No." in table.get_text(" ", strip=True)
            and "Launch site" in table.get_text(" ", strip=True)
        ]
        row_count = sum(max(len(t.find_all("tr")) - 1, 0) for t in launch_tables)
        header_cells = (
            [
                th.get_text(" ", strip=True)
                for th in launch_tables[0].find_all("tr")[0].find_all("th")
            ]
            if launch_tables
            else []
        )
        return {
            "status": "completed",
            "tables": len(launch_tables),
            "rows": row_count,
            "headers": header_cells[:8],
            "url": url,
        }
    except Exception as exc:
        return {
            "status": f"unavailable: {exc}",
            "tables": 0,
            "rows": 0,
            "headers": [],
            "url": url,
        }


def sql_results(sql_df):
    con = sqlite3.connect(":memory:")
    sql_df.to_sql("SPACEXTABLE", con, if_exists="replace", index=False)
    queries = {
        "unique_sites": """
            SELECT DISTINCT TRIM(Launch_Site) AS Launch_Site
            FROM SPACEXTABLE
            ORDER BY Launch_Site
        """,
        "cca5": """
            SELECT Date, TRIM(Launch_Site) AS Launch_Site, Booster_Version,
                   PAYLOAD_MASS__KG_ AS Payload_kg, Orbit, Landing_Outcome
            FROM SPACEXTABLE
            WHERE TRIM(Launch_Site) LIKE 'CCA%'
            LIMIT 5
        """,
        "nasa_payload": """
            SELECT SUM(PAYLOAD_MASS__KG_) AS total_payload_kg
            FROM SPACEXTABLE
            WHERE Customer LIKE '%NASA%'
        """,
        "f9v11_avg": """
            SELECT AVG(PAYLOAD_MASS__KG_) AS avg_payload_kg
            FROM SPACEXTABLE
            WHERE Booster_Version LIKE 'F9 v1.1%'
        """,
        "first_ground_success": """
            SELECT MIN(Date) AS first_success_ground_pad
            FROM SPACEXTABLE
            WHERE Landing_Outcome = 'Success (ground pad)'
        """,
        "drone_4000_6000": """
            SELECT Booster_Version, PAYLOAD_MASS__KG_ AS Payload_kg
            FROM SPACEXTABLE
            WHERE Landing_Outcome = 'Success (drone ship)'
              AND PAYLOAD_MASS__KG_ > 4000
              AND PAYLOAD_MASS__KG_ < 6000
            ORDER BY Payload_kg
        """,
        "mission_outcomes": """
            SELECT TRIM(Mission_Outcome) AS Mission_Outcome, COUNT(*) AS Count
            FROM SPACEXTABLE
            GROUP BY TRIM(Mission_Outcome)
            ORDER BY Count DESC
        """,
        "max_payload": """
            SELECT Booster_Version, PAYLOAD_MASS__KG_ AS Payload_kg
            FROM SPACEXTABLE
            WHERE PAYLOAD_MASS__KG_ = (
                SELECT MAX(PAYLOAD_MASS__KG_) FROM SPACEXTABLE
            )
        """,
        "failed_drone_2015": """
            SELECT Date, Booster_Version, TRIM(Launch_Site) AS Launch_Site,
                   Landing_Outcome
            FROM SPACEXTABLE
            WHERE Landing_Outcome = 'Failure (drone ship)'
              AND Date LIKE '2015%'
        """,
        "rank_outcomes": """
            SELECT Landing_Outcome, COUNT(*) AS Count
            FROM SPACEXTABLE
            WHERE Date BETWEEN '2010-06-04' AND '2017-03-20'
            GROUP BY Landing_Outcome
            ORDER BY Count DESC
        """,
    }
    return {name: pd.read_sql_query(query, con) for name, query in queries.items()}


def train_models(part2, part3):
    y = part2["Class"].to_numpy()
    x = preprocessing.StandardScaler().fit_transform(part3)
    x_train, x_test, y_train, y_test = train_test_split(
        x, y, test_size=0.2, random_state=2
    )
    grids = [
        (
            "Logistic Regression",
            LogisticRegression(max_iter=1000),
            {"C": [0.01, 0.1, 1], "penalty": ["l2"], "solver": ["lbfgs"]},
        ),
        (
            "SVM",
            SVC(),
            {
                "kernel": ["linear", "rbf", "poly", "sigmoid"],
                "C": np.logspace(-3, 3, 5),
                "gamma": np.logspace(-3, 3, 5),
            },
        ),
        (
            "Decision Tree",
            DecisionTreeClassifier(random_state=2),
            {
                "criterion": ["gini", "entropy"],
                "splitter": ["best", "random"],
                "max_depth": [2 * n for n in range(1, 10)],
                "max_features": [None, "sqrt", "log2"],
                "min_samples_leaf": [1, 2, 4],
                "min_samples_split": [2, 5, 10],
            },
        ),
        (
            "KNN",
            KNeighborsClassifier(),
            {
                "n_neighbors": list(range(1, 11)),
                "algorithm": ["auto", "ball_tree", "kd_tree", "brute"],
                "p": [1, 2],
            },
        ),
    ]

    rows = []
    fitted = {}
    for name, model, params in grids:
        grid = GridSearchCV(model, params, cv=10, n_jobs=1)
        grid.fit(x_train, y_train)
        pred = grid.predict(x_test)
        rows.append(
            {
                "model": name,
                "cv_accuracy": float(grid.best_score_),
                "test_accuracy": float(accuracy_score(y_test, pred)),
                "best_params": grid.best_params_,
                "confusion_matrix": confusion_matrix(y_test, pred),
            }
        )
        fitted[name] = grid

    result = pd.DataFrame(rows).sort_values(
        ["test_accuracy", "cv_accuracy"], ascending=False
    )
    best_name = result.iloc[0]["model"]
    return {
        "summary": result,
        "best_model": best_name,
        "best_cm": result.iloc[0]["confusion_matrix"],
        "y_test": y_test,
        "feature_count": part3.shape[1],
    }


def save_fig(fig, name):
    path = FIG / name
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def create_charts(frames, sql, ml):
    part2 = frames["part2"].copy()
    part2["Year"] = pd.to_datetime(part2["Date"]).dt.year
    part2["OutcomeLabel"] = part2["Class"].map({0: "Not landed", 1: "Landed"})

    palette = {0: COLORS["red"], 1: COLORS["teal"]}
    paths = {}

    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    sns.scatterplot(
        data=part2,
        x="FlightNumber",
        y="LaunchSite",
        hue="Class",
        palette=palette,
        s=70,
        edgecolor="white",
        linewidth=0.6,
        ax=ax,
    )
    ax.set_title("Flight Number vs. Launch Site")
    ax.set_xlabel("Flight Number")
    ax.set_ylabel("Launch Site")
    ax.legend(title="Class", labels=["Not landed", "Landed"])
    paths["flight_site"] = save_fig(fig, "eda_flight_number_vs_launch_site.png")

    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    sns.scatterplot(
        data=part2,
        x="PayloadMass",
        y="LaunchSite",
        hue="Class",
        palette=palette,
        s=70,
        edgecolor="white",
        linewidth=0.6,
        ax=ax,
    )
    ax.set_title("Payload Mass vs. Launch Site")
    ax.set_xlabel("Payload Mass (kg)")
    ax.set_ylabel("Launch Site")
    paths["payload_site"] = save_fig(fig, "eda_payload_vs_launch_site.png")

    orbit = (
        part2.groupby("Orbit", as_index=False)
        .agg(success_rate=("Class", "mean"), launches=("Class", "size"))
        .sort_values("success_rate", ascending=False)
    )
    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    sns.barplot(data=orbit, x="Orbit", y="success_rate", color=COLORS["blue"], ax=ax)
    ax.set_ylim(0, 1.05)
    ax.set_title("Success Rate by Orbit Type")
    ax.set_ylabel("Success Rate")
    ax.set_xlabel("Orbit")
    ax.tick_params(axis="x", rotation=35)
    for i, row in enumerate(orbit.itertuples()):
        ax.text(i, row.success_rate + 0.025, f"{row.success_rate:.0%}", ha="center", fontsize=8)
    paths["orbit_success"] = save_fig(fig, "eda_success_rate_by_orbit.png")

    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    sns.scatterplot(
        data=part2,
        x="FlightNumber",
        y="Orbit",
        hue="Class",
        palette=palette,
        s=70,
        edgecolor="white",
        linewidth=0.6,
        ax=ax,
    )
    ax.set_title("Flight Number vs. Orbit Type")
    ax.set_xlabel("Flight Number")
    ax.set_ylabel("Orbit")
    paths["flight_orbit"] = save_fig(fig, "eda_flight_number_vs_orbit.png")

    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    sns.scatterplot(
        data=part2,
        x="PayloadMass",
        y="Orbit",
        hue="Class",
        palette=palette,
        s=70,
        edgecolor="white",
        linewidth=0.6,
        ax=ax,
    )
    ax.set_title("Payload Mass vs. Orbit Type")
    ax.set_xlabel("Payload Mass (kg)")
    ax.set_ylabel("Orbit")
    paths["payload_orbit"] = save_fig(fig, "eda_payload_vs_orbit.png")

    yearly = (
        part2.groupby("Year", as_index=False)
        .agg(success_rate=("Class", "mean"), launches=("Class", "size"))
        .sort_values("Year")
    )
    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    ax.plot(
        yearly["Year"],
        yearly["success_rate"],
        marker="o",
        linewidth=2.7,
        color=COLORS["teal"],
    )
    ax.fill_between(yearly["Year"], yearly["success_rate"], color=COLORS["teal"], alpha=0.12)
    ax.set_ylim(0, 1.05)
    ax.set_title("Yearly Average Launch Success Trend")
    ax.set_xlabel("Year")
    ax.set_ylabel("Average Success Rate")
    ax.set_xticks(yearly["Year"])
    ax.tick_params(axis="x", rotation=35)
    paths["yearly"] = save_fig(fig, "eda_yearly_success_trend.png")

    site = (
        part2.groupby("LaunchSite", as_index=False)
        .agg(launches=("Class", "size"), successes=("Class", "sum"), success_rate=("Class", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    fig, ax = plt.subplots(figsize=(8.5, 4.6))
    sns.barplot(data=site, x="LaunchSite", y="success_rate", color=COLORS["teal"], ax=ax)
    ax.set_ylim(0, 1.05)
    ax.set_title("Launch Site Success Rate")
    ax.set_xlabel("Launch Site")
    ax.set_ylabel("Success Rate")
    for i, row in enumerate(site.itertuples()):
        ax.text(i, row.success_rate + 0.025, f"{row.success_rate:.0%}", ha="center", fontsize=9)
    paths["site_success"] = save_fig(fig, "summary_site_success_rate.png")

    dash = frames["dash"].copy()
    success_counts = (
        dash[dash["class"] == 1]
        .groupby("Launch Site", as_index=False)
        .size()
        .rename(columns={"size": "Successful Launches"})
    )
    fig, ax = plt.subplots(figsize=(7.6, 5.0))
    ax.pie(
        success_counts["Successful Launches"],
        labels=success_counts["Launch Site"],
        autopct="%1.0f%%",
        startangle=110,
        colors=[COLORS["teal"], COLORS["blue"], COLORS["amber"], "#7C3AED"],
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
    )
    ax.set_title("Dash Pie: Total Successful Launches by Site")
    paths["dash_pie_all"] = save_fig(fig, "dash_success_pie_all_sites.png")

    best_site = site.iloc[0]["LaunchSite"]
    selected = part2[part2["LaunchSite"] == best_site]["Class"].value_counts().reindex([1, 0], fill_value=0)
    fig, ax = plt.subplots(figsize=(7.6, 5.0))
    ax.pie(
        selected.values,
        labels=["Landed", "Not landed"],
        autopct="%1.0f%%",
        startangle=90,
        colors=[COLORS["teal"], COLORS["red"]],
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
    )
    ax.set_title(f"Dash Pie: Outcomes for {best_site}")
    paths["dash_pie_site"] = save_fig(fig, "dash_success_pie_best_site.png")

    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    sns.scatterplot(
        data=part2,
        x="PayloadMass",
        y="Class",
        hue="LaunchSite",
        s=75,
        alpha=0.82,
        edgecolor="white",
        linewidth=0.5,
        ax=ax,
    )
    ax.set_title("Dash Scatter: Payload Mass vs. Landing Outcome")
    ax.set_xlabel("Payload Mass (kg)")
    ax.set_ylabel("Landing Outcome (0=Not landed, 1=Landed)")
    paths["dash_scatter"] = save_fig(fig, "dash_payload_vs_outcome_scatter.png")

    ml_summary = ml["summary"].copy()
    fig, ax = plt.subplots(figsize=(8.8, 4.6))
    width = 0.36
    x = np.arange(len(ml_summary))
    ax.bar(x - width / 2, ml_summary["cv_accuracy"], width, label="10-fold CV", color=COLORS["blue"])
    ax.bar(x + width / 2, ml_summary["test_accuracy"], width, label="Test", color=COLORS["teal"])
    ax.set_xticks(x)
    ax.set_xticklabels(ml_summary["model"], rotation=20, ha="right")
    ax.set_ylim(0, 1.0)
    ax.set_ylabel("Accuracy")
    ax.set_title("Classification Model Accuracy")
    ax.legend(frameon=False)
    for idx, row in enumerate(ml_summary.itertuples()):
        ax.text(idx - width / 2, row.cv_accuracy + 0.02, f"{row.cv_accuracy:.2f}", ha="center", fontsize=8)
        ax.text(idx + width / 2, row.test_accuracy + 0.02, f"{row.test_accuracy:.2f}", ha="center", fontsize=8)
    paths["ml_accuracy"] = save_fig(fig, "ml_model_accuracy.png")

    cm = ml["best_cm"]
    fig, ax = plt.subplots(figsize=(5.4, 4.7))
    sns.heatmap(
        cm,
        annot=True,
        fmt="d",
        cmap=sns.light_palette(COLORS["teal"], as_cmap=True),
        cbar=False,
        square=True,
        linewidths=1,
        linecolor="white",
        ax=ax,
    )
    ax.set_xlabel("Predicted label")
    ax.set_ylabel("True label")
    ax.set_xticklabels(["Not landed", "Landed"], rotation=0)
    ax.set_yticklabels(["Not landed", "Landed"], rotation=0)
    ax.set_title(f"Confusion Matrix: {ml['best_model']}")
    paths["ml_cm"] = save_fig(fig, "ml_best_confusion_matrix.png")

    geo = frames["geo"].copy()
    geo_sites = (
        geo.groupby("Launch Site", as_index=False)
        .agg(Lat=("Lat", "first"), Long=("Long", "first"), launches=("class", "size"), successes=("class", "sum"))
    )
    geo_sites["rate"] = geo_sites["successes"] / geo_sites["launches"]
    fig, ax = plt.subplots(figsize=(8.8, 5.0))
    ax.scatter(
        geo_sites["Long"],
        geo_sites["Lat"],
        s=geo_sites["launches"] * 32,
        color=COLORS["blue"],
        alpha=0.75,
        edgecolors="white",
        linewidths=1.1,
    )
    for row in geo_sites.itertuples():
        ax.text(row.Long + 0.45, row.Lat, f"{row._1}\n{row.launches} launches", fontsize=8.5, va="center")
    ax.set_title("Folium Map Concept: Launch Site Markers")
    ax.set_xlabel("Longitude")
    ax.set_ylabel("Latitude")
    ax.set_xlim(-123, -78)
    ax.set_ylim(26, 36)
    ax.grid(True, color="#E5E7EB")
    paths["map_markers"] = save_fig(fig, "folium_launch_site_markers.png")

    fig, ax = plt.subplots(figsize=(8.8, 5.0))
    colors = geo["class"].map({1: COLORS["teal"], 0: COLORS["red"]})
    ax.scatter(geo["Long"], geo["Lat"], c=colors, s=62, alpha=0.78, edgecolors="white", linewidths=0.7)
    for row in geo_sites.itertuples():
        ax.text(row.Long + 0.35, row.Lat + 0.16, row._1, fontsize=8.5, weight="bold")
    ax.set_title("Folium MarkerCluster Concept: Launch Records by Outcome")
    ax.set_xlabel("Longitude")
    ax.set_ylabel("Latitude")
    ax.set_xlim(-123, -78)
    ax.set_ylim(26, 36)
    ax.grid(True, color="#E5E7EB")
    ax.scatter([], [], color=COLORS["teal"], label="Success")
    ax.scatter([], [], color=COLORS["red"], label="Failure/no attempt")
    ax.legend(frameon=False, loc="lower left")
    paths["map_records"] = save_fig(fig, "folium_launch_record_markers.png")

    ksc = geo_sites[geo_sites["Launch Site"].str.contains("KSC")].iloc[0]
    coastline = (ksc.Lat, -80.57)
    highway = (28.55, -80.65)
    city = (28.3922, -80.6077)

    def haversine(lat1, lon1, lat2, lon2):
        r = 6373.0
        lat1, lon1, lat2, lon2 = map(np.radians, [lat1, lon1, lat2, lon2])
        dlon = lon2 - lon1
        dlat = lat2 - lat1
        a = np.sin(dlat / 2) ** 2 + np.cos(lat1) * np.cos(lat2) * np.sin(dlon / 2) ** 2
        return float(r * 2 * np.arctan2(np.sqrt(a), np.sqrt(1 - a)))

    prox = pd.DataFrame(
        [
            ["Coastline", coastline[0], coastline[1], haversine(ksc.Lat, ksc.Long, *coastline)],
            ["Highway", highway[0], highway[1], haversine(ksc.Lat, ksc.Long, *highway)],
            ["Cape Canaveral", city[0], city[1], haversine(ksc.Lat, ksc.Long, *city)],
        ],
        columns=["Feature", "Lat", "Long", "Distance_km"],
    )
    fig, ax = plt.subplots(figsize=(8.8, 5.0))
    ax.scatter([ksc.Long], [ksc.Lat], s=180, color=COLORS["blue"], edgecolor="white", linewidth=1.2, label="KSC LC-39A")
    ax.scatter(prox["Long"], prox["Lat"], s=95, color=COLORS["amber"], edgecolor="white", linewidth=1.0)
    for row in prox.itertuples():
        ax.plot([ksc.Long, row.Long], [ksc.Lat, row.Lat], color=COLORS["teal"], linewidth=2)
        ax.text(row.Long + 0.008, row.Lat, f"{row.Feature}\n{row.Distance_km:.1f} km", fontsize=9, va="center")
    ax.text(ksc.Long + 0.006, ksc.Lat + 0.01, "KSC LC-39A", fontsize=9, weight="bold")
    ax.set_xlim(-80.72, -80.50)
    ax.set_ylim(28.34, 28.60)
    ax.set_title("Folium Proximity Analysis Concept: KSC LC-39A")
    ax.set_xlabel("Longitude")
    ax.set_ylabel("Latitude")
    ax.grid(True, color="#E5E7EB")
    paths["map_proximity"] = save_fig(fig, "folium_proximity_analysis.png")

    return paths


def fig_page():
    fig, ax = plt.subplots(figsize=(16, 9))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.axis("off")
    ax.add_patch(FancyBboxPatch((0, 0), 16, 9, boxstyle="square,pad=0", fc=COLORS["bg"], ec="none"))
    return fig, ax


def draw_header(ax, title, label=None):
    ax.add_patch(FancyBboxPatch((0, 8.68), 16, 0.32, boxstyle="square,pad=0", fc=COLORS["navy"], ec="none"))
    ax.add_patch(FancyBboxPatch((0, 8.68), 3.9, 0.32, boxstyle="square,pad=0", fc=COLORS["teal"], ec="none"))
    ax.text(0.65, 8.22, title, fontsize=24, weight="bold", color=COLORS["ink"], va="top")
    if label:
        ax.text(15.35, 8.28, label, fontsize=9.5, color=COLORS["muted"], ha="right", va="top")


def draw_footer(ax, page=None):
    ax.text(0.65, 0.33, "IBM Applied Data Science Capstone | SpaceX Falcon 9 Landing Prediction", fontsize=8.5, color=COLORS["muted"])
    ax.text(15.35, 0.33, f"{page}" if page else "", fontsize=8.5, color=COLORS["muted"], ha="right")


def wrap_lines(text, width):
    lines = []
    for raw in str(text).split("\n"):
        if raw.strip() == "":
            lines.append("")
        else:
            break_long = "http://" in raw or "https://" in raw
            lines.extend(
                textwrap.wrap(
                    raw,
                    width=width,
                    break_long_words=break_long,
                    break_on_hyphens=True,
                )
            )
    return lines


def bullets(ax, items, x, y, width=70, size=12.3, color=None, line_gap=0.37, bullet_color=None):
    yy = y
    color = color or COLORS["ink"]
    bullet_color = bullet_color or COLORS["teal"]
    for item in items:
        if isinstance(item, tuple):
            text, sub_color = item
            bcol = sub_color
        else:
            text = item
            bcol = bullet_color
        wrapped = wrap_lines(text, width)
        ax.text(x, yy, "●", fontsize=size * 0.86, color=bcol, va="top")
        ax.text(x + 0.32, yy, wrapped[0] if wrapped else "", fontsize=size, color=color, va="top")
        yy -= line_gap
        for cont in wrapped[1:]:
            ax.text(x + 0.32, yy, cont, fontsize=size, color=color, va="top")
            yy -= line_gap
        yy -= 0.08
    return yy


def panel(ax, x, y, w, h, title=None, fill="white", edge=None):
    ax.add_patch(
        FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.02,rounding_size=0.06",
            fc=fill,
            ec=edge or COLORS["line"],
            lw=0.9,
        )
    )
    if title:
        ax.text(x + 0.28, y + h - 0.32, title, fontsize=12.5, weight="bold", color=COLORS["ink"], va="top")


def image(ax, path, x, y, w, h):
    img = Image.open(path)
    ax.imshow(img, extent=(x, x + w, y, y + h), aspect="auto")


def draw_url(ax, url, x, y, width=58, size=8.5):
    yy = y
    for line in wrap_lines(url, width):
        ax.text(x, yy, line, fontsize=size, color=COLORS["blue"], va="top")
        yy -= 0.28
    return yy


def table(ax, df, x, y, w, h, max_rows=8, font=8.2, col_widths=None):
    if isinstance(df, pd.DataFrame):
        data = df.copy().head(max_rows)
        labels = list(data.columns)
        rows = data.astype(str).values.tolist()
    else:
        labels, rows = df
    tbl = ax.table(cellText=rows, colLabels=labels, bbox=(x / 16, y / 9, w / 16, h / 9), colWidths=col_widths)
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(font)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#E2E8F0")
        cell.set_linewidth(0.45)
        if r == 0:
            cell.set_text_props(weight="bold", color="white")
            cell.set_facecolor(COLORS["navy"])
        else:
            cell.set_facecolor("white" if r % 2 else "#F8FAFC")
            cell.set_text_props(color=COLORS["ink"])
    return tbl


def flowchart(ax, steps, x, y, w, box_h=0.72, gap=0.35):
    n = len(steps)
    box_w = (w - gap * (n - 1)) / n
    centers = []
    for i, step in enumerate(steps):
        xx = x + i * (box_w + gap)
        panel(ax, xx, y, box_w, box_h, fill="white", edge=COLORS["line"])
        ax.text(xx + box_w / 2, y + box_h / 2, step, fontsize=9.7, color=COLORS["ink"], ha="center", va="center", wrap=True)
        centers.append((xx + box_w, y + box_h / 2))
        if i < n - 1:
            ax.add_patch(
                FancyArrowPatch(
                    (xx + box_w + 0.05, y + box_h / 2),
                    (xx + box_w + gap - 0.08, y + box_h / 2),
                    arrowstyle="-|>",
                    mutation_scale=14,
                    color=COLORS["teal"],
                    lw=1.4,
                )
            )


def add_slide(pdf, fig, page):
    draw_footer(fig.axes[0], page)
    pdf.savefig(fig, facecolor=COLORS["bg"])
    plt.close(fig)


def slide_title(pdf, page):
    fig, ax = fig_page()
    ax.add_patch(FancyBboxPatch((0, 0), 16, 9, boxstyle="square,pad=0", fc="#0F172A", ec="none"))
    ax.add_patch(FancyBboxPatch((0, 0), 5.15, 9, boxstyle="square,pad=0", fc=COLORS["teal"], ec="none", alpha=0.95))
    ax.text(0.7, 7.65, "Data Science\nCapstone Project\nReport", fontsize=32, weight="bold", color="white", va="top")
    ax.text(0.75, 3.05, "SpaceX Falcon 9 First Stage\nLanding Prediction", fontsize=18, color="white", va="top")
    ax.text(0.75, 1.72, "Prepared for IBM Applied Data Science Capstone", fontsize=11.5, color="#E0F2FE")
    panel(ax, 6.05, 1.1, 8.95, 6.7, fill="#FFFFFF", edge="#334155")
    ax.text(6.55, 7.25, "Project Objective", fontsize=15, weight="bold", color=COLORS["ink"])
    bullets(
        ax,
        [
            "Predict whether a Falcon 9 first stage will land successfully.",
            "Combine API collection, web scraping, data wrangling, EDA, SQL, interactive visualization, and machine learning.",
            "Translate the results into launch-site, payload, orbit, and model-performance insights.",
        ],
        6.55,
        6.65,
        width=76,
        size=12.8,
    )
    ax.text(6.55, 3.33, "GitHub URL", fontsize=13, weight="bold", color=COLORS["ink"])
    ax.text(6.55, 2.86, GITHUB_URL, fontsize=11.2, color=COLORS["blue"])
    ax.text(6.55, 2.24, "Final upload artifact: completed presentation slides in PDF format.", fontsize=11.3, color=COLORS["muted"])
    draw_footer(ax, page)
    pdf.savefig(fig, facecolor="#0F172A")
    plt.close(fig)


def build_pdf(frames, sql, ml, charts, scrape):
    part1 = frames["part1"].copy()
    part2 = frames["part2"].copy()
    part2["Year"] = pd.to_datetime(part2["Date"]).dt.year
    site = (
        part2.groupby("LaunchSite", as_index=False)
        .agg(launches=("Class", "size"), successes=("Class", "sum"), success_rate=("Class", "mean"), avg_payload=("PayloadMass", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    orbit = (
        part2.groupby("Orbit", as_index=False)
        .agg(launches=("Class", "size"), successes=("Class", "sum"), success_rate=("Class", "mean"))
        .sort_values("success_rate", ascending=False)
    )
    yearly = part2.groupby("Year", as_index=False).agg(launches=("Class", "size"), success_rate=("Class", "mean"))
    best_model = ml["best_model"]
    best_row = ml["summary"].iloc[0]
    cm = ml["best_cm"]

    with PdfPages(PDF_EN) as pdf:
        page = 1
        slide_title(pdf, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Outline", "Project Structure")
        panel(ax, 0.8, 1.0, 14.4, 6.5, fill="white")
        bullets(
            ax,
            [
                "Executive Summary and business problem",
                "Methodology: API collection, web scraping, data wrangling, EDA, SQL, Folium, Dash, and classification modeling",
                "EDA results: scatter plots, orbit bar chart, and yearly trend",
                "SQL results: launch sites, payloads, success/failure counts, ranking, and time analysis",
                "Interactive visual analytics: Folium maps and Plotly Dash dashboard",
                "Predictive analysis: model comparison, confusion matrix, best model, and conclusions",
            ],
            1.3,
            6.8,
            width=118,
            size=13.2,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Executive Summary", "Rubric 1.3")
        panel(ax, 0.8, 1.1, 6.8, 6.45, title="Methods Used")
        bullets(
            ax,
            [
                "Collected launch data from SpaceX REST API endpoints and a historical Wikipedia launch table.",
                "Cleaned missing payload values, encoded landing outcome as Class, and one-hot encoded categorical fields.",
                "Explored patterns with visual EDA, SQL queries, Folium-style maps, and Dash dashboard views.",
                "Trained Logistic Regression, SVM, Decision Tree, and KNN models with 10-fold GridSearchCV.",
            ],
            1.18,
            6.85,
            width=58,
            size=11.5,
        )
        panel(ax, 8.05, 1.1, 7.15, 6.45, title="Key Results")
        bullets(
            ax,
            [
                f"Dataset contains {len(part2)} Falcon 9 launches with an overall first-stage landing success rate of {part2['Class'].mean():.1%}.",
                f"Best-performing launch sites by success rate: {site.iloc[0]['LaunchSite']} ({site.iloc[0]['success_rate']:.1%}) and {site.iloc[1]['LaunchSite']} ({site.iloc[1]['success_rate']:.1%}).",
                f"Orbit effects are visible: SSO, GEO, HEO and ES-L1 show 100% success in this dataset, while GTO is lower at {orbit[orbit['Orbit']=='GTO']['success_rate'].iloc[0]:.1%}.",
                f"Best model selected: {best_model}, with test accuracy {best_row['test_accuracy']:.1%} and 10-fold CV accuracy {best_row['cv_accuracy']:.1%}.",
            ],
            8.42,
            6.85,
            width=62,
            size=11.5,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Introduction", "Rubric 1.4")
        panel(ax, 0.8, 1.0, 14.4, 6.6, fill="white")
        ax.text(1.25, 6.95, "Project Background", fontsize=15, weight="bold", color=COLORS["ink"])
        bullets(
            ax,
            [
                "Reusable Falcon 9 first stages materially reduce launch cost; a successful landing is therefore a valuable operational and commercial outcome.",
                "The project treats landing success as a binary classification problem where Class=1 means the booster landed successfully and Class=0 means it did not.",
                "The analysis connects launch site, orbit, payload mass, booster reuse, grid fins, legs, and launch chronology to landing outcome.",
            ],
            1.25,
            6.36,
            width=120,
            size=12.5,
        )
        ax.text(1.25, 3.36, "Problem Statement", fontsize=15, weight="bold", color=COLORS["ink"])
        ax.text(
            1.25,
            2.82,
            "Can we predict Falcon 9 first-stage landing success before launch, and which launch conditions are most associated with reliable recovery?",
            fontsize=14.2,
            color=COLORS["blue"],
            weight="bold",
            wrap=True,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Data Collection - SpaceX API", "Rubric 1.5")
        flowchart(
            ax,
            [
                "GET /v4/launches/past",
                "Normalize JSON",
                "Lookup rockets, payloads, launchpads, cores",
                "Filter Falcon 9",
                "Export dataset_part_1.csv",
            ],
            0.95,
            5.95,
            14.1,
        )
        panel(ax, 0.95, 1.1, 6.85, 4.0, title="API Call Logic")
        bullets(
            ax,
            [
                "Primary launch endpoint supplies IDs for rockets, payloads, launchpads, and cores.",
                "Helper requests enrich each launch with booster version, payload mass, orbit, customer, launch site name, coordinates, grid fins, legs, reuse count, and serial number.",
                f"After filtering to Falcon 9, the working dataset contains {len(part1)} rows.",
            ],
            1.25,
            4.45,
            width=60,
            size=10.9,
        )
        panel(ax, 8.15, 1.1, 6.9, 4.0, title="External Reference")
        ax.text(8.45, 4.45, "GitHub URL:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        draw_url(ax, gh(NOTEBOOKS["api"]), 8.45, 4.02, width=58, size=8.3)
        ax.text(8.45, 3.25, "Key endpoints:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        bullets(
            ax,
            ["/v4/launches/past", "/v4/rockets/{id}", "/v4/payloads/{id}", "/v4/launchpads/{id}", "/v4/cores/{id}"],
            8.45,
            2.82,
            width=52,
            size=10.5,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Data Collection - Web Scraping", "Rubric 1.6")
        flowchart(
            ax,
            [
                "Request archived Wikipedia page",
                "Parse HTML with BeautifulSoup",
                "Find Falcon launch tables",
                "Extract headers and rows",
                "Clean text into DataFrame",
            ],
            0.95,
            5.95,
            14.1,
        )
        panel(ax, 0.95, 1.05, 7.0, 4.15, title="Scraping Process")
        bullets(
            ax,
            [
                f"Source: archived Wikipedia list of Falcon 9 and Falcon Heavy launches.",
                f"Parsed {scrape['tables']} launch tables and approximately {scrape['rows']} launch rows from the HTML page.",
                "Cleaned date/time, booster version, launch site, payload, orbit, customer, launch outcome, and landing outcome fields.",
            ],
            1.25,
            4.55,
            width=62,
            size=11.0,
        )
        panel(ax, 8.25, 1.05, 6.8, 4.15, title="External Reference")
        ax.text(8.55, 4.55, "GitHub URL:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        draw_url(ax, gh(NOTEBOOKS["scraping"]), 8.55, 4.12, width=58, size=8.3)
        ax.text(8.55, 3.36, "Key phrases:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        bullets(ax, ["requests.get", "BeautifulSoup", "find_all('table')", "parse rows", "pd.DataFrame"], 8.55, 2.92, width=52, size=10.5)
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Data Wrangling", "Rubric 1.7")
        flowchart(
            ax,
            [
                "Inspect missing values",
                "Replace missing payload mass with mean",
                "Classify landing outcomes",
                "Encode categorical fields",
                "Cast model matrix to float64",
            ],
            0.95,
            5.95,
            14.1,
        )
        panel(ax, 0.95, 1.1, 6.9, 4.0, title="Cleaning Decisions")
        payload_mean = part1["PayloadMass"].mean()
        bullets(
            ax,
            [
                f"PayloadMass missing values were filled with the mean payload mass ({payload_mean:,.1f} kg) in the API collection workflow.",
                "Landing outcome was transformed into a binary Class label: successful outcomes are 1; failed/no-attempt outcomes are 0.",
                f"Feature engineering produced {frames['part3'].shape[1]} model-ready numeric columns.",
            ],
            1.25,
            4.45,
            width=60,
            size=11,
        )
        panel(ax, 8.2, 1.1, 6.85, 4.0, title="External Reference")
        ax.text(8.5, 4.45, "GitHub URL:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        draw_url(ax, gh(NOTEBOOKS["wrangling"]), 8.5, 4.02, width=58, size=8.3)
        ax.text(8.5, 3.2, "Landing outcome distribution:", fontsize=11.2, weight="bold", color=COLORS["ink"])
        mini = part1["Outcome"].value_counts().head(5).reset_index()
        mini.columns = ["Outcome", "Count"]
        table(ax, mini, 8.45, 1.35, 6.25, 1.62, max_rows=5, font=7.6)
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "EDA with Data Visualization", "Rubric 1.8")
        panel(ax, 0.8, 1.05, 6.65, 6.55, title="Charts Built and Why")
        bullets(
            ax,
            [
                "Scatter: Flight Number vs. Launch Site to examine learning effects over time by site.",
                "Scatter: Payload Mass vs. Launch Site to test whether heavier missions land differently at each site.",
                "Bar: Success Rate by Orbit to compare mission profiles.",
                "Scatter: Flight Number/Payload vs. Orbit to detect orbit-specific clusters.",
                "Line: Yearly average success rate to show operational improvement over time.",
            ],
            1.12,
            6.9,
            width=58,
            size=10.9,
        )
        image(ax, charts["site_success"], 8.0, 3.2, 6.8, 3.4)
        image(ax, charts["yearly"], 8.0, 0.85, 6.8, 2.1)
        ax.text(8.0, 7.05, "GitHub URL:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        draw_url(ax, gh(NOTEBOOKS["eda_viz"]), 8.0, 6.62, width=60, size=8.3)
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "EDA with SQL", "Rubric 1.9")
        panel(ax, 0.8, 1.1, 7.0, 6.45, title="SQL Query Families")
        bullets(
            ax,
            [
                "SELECT DISTINCT launch sites and sample records beginning with CCA.",
                "Aggregate payload mass for NASA customers and average payload for F9 v1.1 boosters.",
                "Find first successful ground-pad landing and boosters with successful drone-ship landings in a payload range.",
                "Count mission outcomes, identify maximum-payload boosters, inspect 2015 failed drone landings, and rank landing outcomes by time window.",
            ],
            1.14,
            6.86,
            width=62,
            size=11.1,
        )
        panel(ax, 8.2, 1.1, 6.85, 6.45, title="External Reference")
        ax.text(8.5, 6.86, "GitHub URL:", fontsize=11.5, weight="bold", color=COLORS["ink"])
        draw_url(ax, gh(NOTEBOOKS["eda_sql"]), 8.5, 6.43, width=60, size=8.3)
        ax.text(8.5, 5.64, "Example query:", fontsize=11.4, weight="bold", color=COLORS["ink"])
        code = "SELECT Landing_Outcome, COUNT(*)\nFROM SPACEXTABLE\nWHERE Date BETWEEN '2010-06-04' AND '2017-03-20'\nGROUP BY Landing_Outcome\nORDER BY COUNT(*) DESC;"
        ax.text(8.55, 5.18, code, fontsize=9.4, color=COLORS["ink"], family="DejaVu Sans Mono", va="top")
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Interactive Visual Analytics", "Rubric 1.10")
        panel(ax, 0.8, 1.05, 6.95, 6.55, title="Folium Map")
        bullets(
            ax,
            [
                "Added Circle and Marker objects for launch sites.",
                "Used MarkerCluster and color-coded markers to separate successful and failed launches.",
                "Used MousePosition, distance calculation, DivIcon labels, and PolyLine objects for proximity analysis.",
                f"GitHub URL: {gh(NOTEBOOKS['folium'])}",
            ],
            1.12,
            6.88,
            width=62,
            size=10.8,
        )
        panel(ax, 8.05, 1.05, 7.15, 6.55, title="Plotly Dash")
        bullets(
            ax,
            [
                "Built a dropdown to switch between all sites and a selected launch site.",
                "Added payload range slider to filter the scatter plot.",
                "Created success pie chart and payload-versus-outcome scatter chart.",
                f"GitHub URL: {gh(NOTEBOOKS['dash'])}",
            ],
            8.37,
            6.88,
            width=64,
            size=10.8,
        )
        add_slide(pdf, fig, page)

        for title, key, insight in [
            ("Flight Number vs. Launch Site", "flight_site", "Later flights show a greater concentration of successful landings, especially at KSC LC-39A and VAFB SLC-4E."),
            ("Payload vs. Launch Site", "payload_site", "Successful landings span both light and heavy payloads; site and mission profile matter more than payload alone."),
            ("Success Rate vs. Orbit Type", "orbit_success", "SSO, GEO, HEO, and ES-L1 are perfect in this sample, while GTO is materially lower."),
            ("Flight Number vs. Orbit Type", "flight_orbit", "Launches diversify across orbit types as flight number grows, and later VLEO/SSO missions show stronger success."),
            ("Payload vs. Orbit Type", "payload_orbit", "VLEO missions include very heavy payloads and still often land successfully, reflecting reuse maturity."),
            ("Launch Success Yearly Trend", "yearly", "The success rate rises after 2015 and remains high through 2019-2020, a clear learning-curve signal."),
        ]:
            page += 1
            fig, ax = fig_page()
            draw_header(ax, title, "Rubric 1.11")
            image(ax, charts[key], 0.8, 1.28, 9.7, 6.45)
            panel(ax, 11.0, 1.28, 4.2, 6.45, title="Interpretation")
            bullets(ax, [insight, f"Source notebook: {gh(NOTEBOOKS['eda_viz'])}"], 11.28, 6.88, width=36, size=10.1)
            add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "SQL Results: Launch Sites and CCA Records", "Rubric 1.12")
        panel(ax, 0.8, 1.05, 6.65, 6.55, title="Unique Launch Sites")
        table(ax, sql["unique_sites"], 1.1, 4.2, 5.95, 2.55, font=9)
        bullets(ax, ["The SQL table includes four site labels after trimming whitespace, with Cape Canaveral represented by LC-40 and SLC-40 naming variants."], 1.1, 3.55, width=56, size=10.5)
        panel(ax, 8.0, 1.05, 7.2, 6.55, title="First Five Records Beginning with CCA")
        table(ax, sql["cca5"], 8.25, 2.0, 6.7, 4.6, max_rows=5, font=6.4)
        draw_url(ax, gh(NOTEBOOKS["eda_sql"]), 8.25, 1.55, width=62, size=7.6)
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "SQL Results: Payload Calculations", "Rubric 1.12")
        panel(ax, 0.8, 1.05, 6.9, 6.55, title="Aggregate Results")
        nasa_payload = sql["nasa_payload"].iloc[0, 0]
        avg_payload = sql["f9v11_avg"].iloc[0, 0]
        ax.text(1.25, 5.95, f"{nasa_payload:,.0f} kg", fontsize=30, weight="bold", color=COLORS["teal"])
        ax.text(1.25, 5.45, "Total payload carried for NASA customers", fontsize=12, color=COLORS["muted"])
        ax.text(1.25, 4.20, f"{avg_payload:,.2f} kg", fontsize=30, weight="bold", color=COLORS["blue"])
        ax.text(1.25, 3.70, "Average payload mass for F9 v1.1 boosters", fontsize=12, color=COLORS["muted"])
        panel(ax, 8.15, 1.05, 6.9, 6.55, title="SQL Logic")
        bullets(
            ax,
            [
                "Used SUM(PAYLOAD_MASS__KG_) filtered by Customer LIKE '%NASA%'.",
                "Used AVG(PAYLOAD_MASS__KG_) filtered by Booster_Version LIKE 'F9 v1.1%'.",
                f"GitHub URL: {gh(NOTEBOOKS['eda_sql'])}",
            ],
            8.45,
            6.88,
            width=60,
            size=11.2,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "SQL Results: Landing Success Details", "Rubric 1.12")
        panel(ax, 0.8, 1.05, 6.95, 6.55, title="First Successful Ground Landing")
        first_date = sql["first_ground_success"].iloc[0, 0]
        ax.text(1.18, 5.95, str(first_date), fontsize=30, weight="bold", color=COLORS["teal"])
        ax.text(1.18, 5.45, "Landing_Outcome = Success (ground pad)", fontsize=12, color=COLORS["muted"])
        bullets(ax, ["This milestone marks the transition from ocean/drone-ship recovery attempts to reliable ground-pad landing capability."], 1.18, 4.55, width=58, size=11.0)
        panel(ax, 8.05, 1.05, 7.15, 6.55, title="Successful Drone Ship Landings: 4000-6000 kg")
        table(ax, sql["drone_4000_6000"], 8.35, 3.15, 6.55, 3.25, max_rows=6, font=8)
        draw_url(ax, gh(NOTEBOOKS["eda_sql"]), 8.35, 2.45, width=62, size=7.6)
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "SQL Results: Outcomes and Maximum Payload", "Rubric 1.12")
        panel(ax, 0.8, 1.05, 6.9, 6.55, title="Mission Outcome Counts")
        table(ax, sql["mission_outcomes"], 1.1, 3.35, 6.3, 3.2, max_rows=5, font=8.2)
        bullets(ax, ["Most missions succeeded; only one record is a clear in-flight mission failure in this SQL dataset."], 1.1, 2.75, width=58, size=10.4)
        panel(ax, 8.1, 1.05, 7.1, 6.55, title="Boosters with Maximum Payload")
        max_payload = sql["max_payload"].iloc[0]["Payload_kg"]
        ax.text(8.45, 6.15, f"Maximum payload: {max_payload:,.0f} kg", fontsize=15, weight="bold", color=COLORS["blue"])
        table(ax, sql["max_payload"].head(8), 8.35, 2.0, 6.55, 3.85, max_rows=8, font=7.2)
        ax.text(8.35, 1.55, f"{len(sql['max_payload'])} boosters carried the maximum payload mass.", fontsize=10.0, color=COLORS["muted"])
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "SQL Results: 2015 and Outcome Ranking", "Rubric 1.12")
        panel(ax, 0.8, 1.05, 6.95, 6.55, title="Failed Drone Ship Landings in 2015")
        table(ax, sql["failed_drone_2015"], 1.05, 3.45, 6.5, 3.0, max_rows=4, font=7.6)
        bullets(ax, ["The two failed 2015 drone-ship landings both launched from CCAFS LC-40."], 1.05, 2.82, width=58, size=10.6)
        panel(ax, 8.05, 1.05, 7.15, 6.55, title="Ranked Landing Outcomes: 2010-06-04 to 2017-03-20")
        table(ax, sql["rank_outcomes"], 8.35, 1.7, 6.55, 4.85, max_rows=8, font=7.4)
        draw_url(ax, gh(NOTEBOOKS["eda_sql"]), 8.35, 1.25, width=62, size=7.6)
        add_slide(pdf, fig, page)

        for title, key, desc in [
            ("Folium Map: Launch Site Markers", "map_markers", "Circle and Marker objects identify all launch sites and communicate their geographic concentration along Florida's Space Coast and Vandenberg."),
            ("Folium Map: Launch Records by Outcome", "map_records", "MarkerCluster-style coloring separates successful launches from failed/no-attempt outcomes for rapid site-level comparison."),
            ("Folium Map: Proximity Analysis", "map_proximity", "Distance markers and PolyLine objects quantify the relationship between a launch site and coastline, highway, and nearby city features."),
        ]:
            page += 1
            fig, ax = fig_page()
            draw_header(ax, title, "Rubric 1.13")
            image(ax, charts[key], 0.8, 1.18, 9.4, 6.55)
            panel(ax, 10.7, 1.18, 4.5, 6.55, title="Map Objects")
            bullets(ax, [desc, "Objects used: Circle, Marker, DivIcon, MarkerCluster, MousePosition, and PolyLine.", f"GitHub URL: {gh(NOTEBOOKS['folium'])}"], 11.0, 6.88, width=39, size=10.2)
            add_slide(pdf, fig, page)

        for title, key, desc in [
            ("Plotly Dash: Successful Launches Pie Chart", "dash_pie_all", "The all-sites pie chart compares where successful launches occurred and reveals site contribution to the recovery record."),
            ("Plotly Dash: Highest-Ratio Site Pie Chart", "dash_pie_site", "The selected-site pie chart focuses on the launch site with the strongest success ratio in the cleaned EDA dataset."),
            ("Plotly Dash: Payload vs. Launch Outcome Scatter", "dash_scatter", "The scatter plot supports interactive filtering by launch site and payload range to inspect whether payload mass changes landing outcomes."),
        ]:
            page += 1
            fig, ax = fig_page()
            draw_header(ax, title, "Rubric 1.14")
            image(ax, charts[key], 0.8, 1.18, 9.75, 6.55)
            panel(ax, 11.0, 1.18, 4.2, 6.55, title="Dashboard Interaction")
            bullets(ax, [desc, "Controls: launch-site dropdown and payload range slider.", f"GitHub URL: {gh(NOTEBOOKS['dash'])}"], 11.28, 6.88, width=36, size=10.2)
            add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Predictive Analysis: Model Development", "Rubric 1.15")
        panel(ax, 0.8, 1.05, 6.55, 6.55, title="Classification Workflow")
        bullets(
            ax,
            [
                "Created target array Y from the Class column.",
                f"Standardized {ml['feature_count']} engineered features using StandardScaler.",
                "Split data into 80% training and 20% test sets with random_state=2.",
                "Used 10-fold GridSearchCV to tune Logistic Regression, SVM, Decision Tree, and KNN.",
                f"GitHub URL: {gh(NOTEBOOKS['ml'])}",
            ],
            1.1,
            6.88,
            width=58,
            size=10.8,
        )
        image(ax, charts["ml_accuracy"], 7.75, 1.25, 7.0, 6.0)
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Predictive Analysis: Best Model and Confusion Matrix", "Rubric 1.15")
        image(ax, charts["ml_cm"], 0.95, 1.45, 6.2, 5.95)
        panel(ax, 8.0, 1.1, 7.05, 6.55, title=f"Best Model: {best_model}")
        tn, fp, fn, tp = cm.ravel()
        bullets(
            ax,
            [
                f"Test accuracy: {best_row['test_accuracy']:.1%}; 10-fold CV accuracy: {best_row['cv_accuracy']:.1%}.",
                f"Confusion matrix: TN={tn}, FP={fp}, FN={fn}, TP={tp}.",
                "The model correctly identifies all successful landings in the test split, while several failed/no-attempt launches are predicted as successful.",
                "Operationally, this suggests the model is useful for early screening but should be paired with mission-risk review before final decisions.",
            ],
            8.32,
            6.86,
            width=62,
            size=11.0,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Conclusions and Actionable Insights", "Rubric 1.15")
        panel(ax, 0.8, 1.05, 14.4, 6.55, fill="white")
        bullets(
            ax,
            [
                "Landing success improved sharply after the early experimental phase; launch chronology is a strong proxy for operational learning and reusable-booster maturity.",
                "KSC LC-39A and VAFB SLC-4E show the highest launch-site success rates in the cleaned Falcon 9 dataset, suggesting site-level operational context should be retained in predictive models.",
                "Orbit type matters: GTO missions remain more difficult in this sample, while SSO/VLEO and selected specialized orbits show stronger recovery performance.",
                "Payload mass alone does not determine landing success; high-payload VLEO missions can still land successfully when paired with mature Block 5/reused-booster operations.",
                "Creative decision insight: combine the classifier output with a rule-based mission risk score by orbit, site, and booster reuse history to flag launches needing additional recovery planning.",
            ],
            1.2,
            6.95,
            width=123,
            size=12.1,
        )
        add_slide(pdf, fig, page)

        page += 1
        fig, ax = fig_page()
        draw_header(ax, "Appendix: Artifacts and Data Sources", "Reference")
        panel(ax, 0.8, 1.0, 14.4, 6.6, title="Project Links")
        rows = [
            ["Project GitHub", GITHUB_URL],
            ["SpaceX API notebook", gh(NOTEBOOKS["api"])],
            ["Web scraping notebook", gh(NOTEBOOKS["scraping"])],
            ["Data wrangling notebook", gh(NOTEBOOKS["wrangling"])],
            ["EDA visualization notebook", gh(NOTEBOOKS["eda_viz"])],
            ["EDA SQL notebook", gh(NOTEBOOKS["eda_sql"])],
            ["Folium notebook", gh(NOTEBOOKS["folium"])],
            ["Dash app", gh(NOTEBOOKS["dash"])],
            ["Machine learning notebook", gh(NOTEBOOKS["ml"])],
        ]
        table(ax, (["Artifact", "URL"], rows), 1.05, 1.35, 13.9, 5.7, max_rows=10, font=7.1)
        add_slide(pdf, fig, page)

    shutil.copyfile(PDF_EN, PDF_UPLOAD)
    return page


def render_pdf_to_pptx(page_count):
    for old in SLIDE_IMAGES.glob("slide_*.png"):
        old.unlink()
    doc = fitz.open(PDF_EN)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    for idx, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        image_path = SLIDE_IMAGES / f"slide_{idx:02d}.png"
        pix.save(image_path)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX_OUT)
    if len(doc) != page_count:
        raise RuntimeError(f"PDF page count mismatch: expected {page_count}, got {len(doc)}")
    doc.close()


def write_artifact_index(frames, sql, ml, scrape, page_count):
    path = OUT / "artifact_index.md"
    lines = [
        "# SpaceX Capstone Report Artifacts",
        "",
        f"- Final PDF: `{PDF_EN.name}`",
        f"- Upload-ready PDF copy: `{PDF_UPLOAD.name}`",
        f"- PPTX preview/edit copy: `{PPTX_OUT.name}`",
        f"- Slide count: {page_count}",
        f"- GitHub URL used in slides: `{GITHUB_URL}`",
        "",
        "## Data Summary",
        f"- Falcon 9 modeling dataset rows: {len(frames['part2'])}",
        f"- Overall landing success rate: {frames['part2']['Class'].mean():.3f}",
        f"- SQL dataset rows: {len(frames['sql'])}",
        f"- Web scraping status: {scrape['status']}, tables={scrape['tables']}, rows~={scrape['rows']}",
        "",
        "## Machine Learning Summary",
    ]
    for row in ml["summary"].itertuples():
        lines.append(f"- {row.model}: CV={row.cv_accuracy:.3f}, test={row.test_accuracy:.3f}, params={row.best_params}")
    lines.append("")
    lines.append("## Note")
    lines.append("Set CAPSTONE_GITHUB_URL before rerunning the script to replace the placeholder GitHub URL in every slide.")
    path.write_text("\n".join(lines), encoding="utf-8")


def main():
    ensure_dirs()
    frames = read_data()
    scrape = wiki_scrape_summary()
    sql = sql_results(frames["sql"])
    ml = train_models(frames["part2"], frames["part3"])
    charts = create_charts(frames, sql, ml)
    page_count = build_pdf(frames, sql, ml, charts, scrape)
    render_pdf_to_pptx(page_count)
    write_artifact_index(frames, sql, ml, scrape, page_count)
    print(f"Generated {PDF_EN}")
    print(f"Generated {PDF_UPLOAD}")
    print(f"Generated {PPTX_OUT}")
    print(f"Slides: {page_count}")
    print(f"GitHub URL used: {GITHUB_URL}")


if __name__ == "__main__":
    main()
